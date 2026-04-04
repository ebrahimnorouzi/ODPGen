#!/usr/bin/env python3
"""
ODPGen Presentation Generator
------------------------------
Reads live data from results/summary.csv and the repo, then writes:
  - ODPGen_Presentation.html
  - ODPGen_Presentation.pptx

Run:
    python make_pptx.py             # both files
    python make_pptx.py --html-only
    python make_pptx.py --pptx-only
"""
import argparse, json, re, textwrap
from pathlib import Path
import pandas as pd

# ── Repo roots (relative to this script) ────────────────────────────────────
ROOT       = Path(__file__).parent
SUMMARY    = ROOT / "results" / "summary.csv"
SCENARIOS  = ROOT / "data" / "scenarios" / "pattern_scenarios.json"
CUTOFFS    = ROOT / "data" / "model_cutoffs.csv"
OUTPUTS    = ROOT / "outputs"
HTML_OUT   = ROOT / "ODPGen_Presentation.html"
PPTX_OUT   = ROOT / "ODPGen_Presentation.pptx"

# ── Per-model display config (auto-extended for unknown models) ──────────────
MODEL_META = {
    "gemini-3.1-pro-preview":            dict(name="Gemini 3.1 Pro",   params="Proprietary", provider="Google",    cutoff="~2025", chip="chip-blue",   pptx_color=(0x8a,0xb4,0xf8)),
    "gpt-5.4":                           dict(name="GPT-5.4",          params="Proprietary", provider="OpenAI",    cutoff="~2024", chip="chip-purple", pptx_color=(0xb3,0x9d,0xdb)),
    "meta-llama_Llama-3.1-8B-Instruct":  dict(name="Llama-3.1-8B",    params="8B",          provider="Meta",      cutoff="Dec 2023", chip="chip-green",  pptx_color=(0x4c,0xaf,0x7d)),
    "meta-llama_Llama-2-70b-chat-hf":    dict(name="Llama-2-70B",     params="70B",         provider="Meta",      cutoff="Sep 2022", chip="chip-yellow", pptx_color=(0xf0,0xc0,0x60)),
    "mistralai_Mistral-7B-Instruct-v0.3":dict(name="Mistral-7B",      params="7B",          provider="Mistral AI",cutoff="~2023",    chip="chip-yellow", pptx_color=(0xf0,0xc0,0x60)),
    "bigscience_bloomz-7b1":             dict(name="BLOOMZ-7B1",       params="7B",          provider="BigScience",cutoff="Jul 2022", chip="chip-red",    pptx_color=(0xe0,0x6c,0x6c)),
    "bigscience_bloom-3b":               dict(name="BLOOM-3B",         params="3B",          provider="BigScience",cutoff="Jul 2022", chip="chip-red",    pptx_color=(0xe0,0x6c,0x6c)),
}

SCENARIO_DOMAINS = {
    "2023-133-01": "Causal Bayesian Networks (Sprinkler)",
    "2023-133-02": "Causal Bayesian Networks (Asthma)",
    "2023-134-01": "Solid Data Governance — Data Request",
    "2023-134-02": "Solid Data Governance — Data Log",
    "2023-134-03": "Solid Data Governance — Data Registry",
    "2023-135-01": "Role-Dependent Names",
    "2024-145-01": "Legal Information (Indian Courts)",
    "2025-147-01": "Maritime Ships (AISHIP)",
    "2025-149-01": "DB Access Control (ARGOS)",
    "2025-150-01": "IoT / MQTT Protocol",
    "2025-151-01": "Access Control — Eval. Request",
    "2025-151-02": "State-of-World (ODRL Policy)",
    "2025-153-01": "Temporal Indirection (Finance)",
    "2026-155-01": "Cultural Heritage Survey",
}

QUALITY_COLS = [
    "cq_coverage_ratio", "scenario_coverage_ratio", "hallucination_ratio",
    "class_label_ratio", "property_label_ratio",
    "class_name_f1", "property_name_f1", "combined_f1",
    "class_count", "object_property_count",
]

CONFIG_ORDER = [
    "cq-only", "scenario-cq", "scenario-cq-constraints",
    "scenario-cq-reasoning", "scenario-only",
]

# ═══════════════════════════════════════════════════════════════════════════════
# DATA LOADING
# ═══════════════════════════════════════════════════════════════════════════════

def load_stats() -> dict:
    """Read summary.csv + repo files; return a flat stats dict used by both renderers."""
    df = pd.read_csv(SUMMARY)
    avail_q = [c for c in QUALITY_COLS if c in df.columns]
    success  = df[df["parse_success"] == True]

    # ── per-model tables ────────────────────────────────────────────────────
    ps = df.groupby("model")["parse_success"].agg(["sum","count","mean"])
    ps.columns = ["parsed","total","rate"]
    qual = success.groupby("model")[avail_q].mean() if not success.empty else pd.DataFrame()

    # ── per-config table ────────────────────────────────────────────────────
    cfg_qual = success.groupby("config")[avail_q].mean() if not success.empty else pd.DataFrame()

    # ── model x config heatmap ──────────────────────────────────────────────
    cross = (df.groupby(["model","config"])["parse_success"]
               .agg(["sum","count","mean"])
               .rename(columns={"mean":"rate"}))

    # ── ordered model list (by parse rate desc, then name) ──────────────────
    models_ordered = ps.sort_values("rate", ascending=False).index.tolist()

    # ── configs present (keep canonical order) ──────────────────────────────
    configs_present = [c for c in CONFIG_ORDER if c in df["config"].unique()]

    # ── scenario metadata ────────────────────────────────────────────────────
    scenarios = []
    if SCENARIOS.exists():
        raw = json.loads(SCENARIOS.read_text())
        items = raw if isinstance(raw, list) else raw.get("scenarios", [])
        for s in items:
            sid = s.get("scenario_id","")
            scenarios.append({
                "id":    sid,
                "cqs":   len(s.get("cq_list", [])),
                "year":  int(sid.split("-")[0]) if sid and sid.split("-")[0].isdigit() else 0,
                "domain": SCENARIO_DOMAINS.get(sid, "—"),
                "cq_list": s.get("cq_list", []),
                "scenario_text": s.get("scenario_text","")[:400],
            })

    # ── model cutoffs ────────────────────────────────────────────────────────
    cutoffs = {}
    if CUTOFFS.exists():
        cdf = pd.read_csv(CUTOFFS)
        for _, row in cdf.iterrows():
            cutoffs[str(row.iloc[0])] = str(row.iloc[1])

    # ── functional evaluation ranking (pilot) ────────────────────────────────
    func_rank = []
    _FUNC_CSV = ROOT / "eval" / "functional_eval_ranking.csv"
    if _FUNC_CSV.exists():
        fdf = pd.read_csv(_FUNC_CSV)
        for _, row in fdf.iterrows():
            try:
                func_rank.append(dict(
                    model=str(row["model"]),
                    config=str(row["config"]),
                    avg_numb=float(row.get("avg_numb", 0)),
                    avg_manual_score=float(row.get("avg_manual_score", 0)),
                    avg_cq_pass_rate=float(row.get("avg_cq_pass_rate", 0)),
                    avg_functional_score=float(row.get("avg_functional_score", 0)),
                ))
            except Exception:
                pass

    # ── sample raw responses (one per model, scenario-cq / 2023-133-01) ─────
    SAMPLE_SID  = "2023-133-01"
    SAMPLE_CFG  = "scenario-cq"
    sample_outputs = {}
    for model in models_ordered:
        for cfg in [SAMPLE_CFG, "scenario-only", "cq-only"]:
            p = OUTPUTS / model / cfg / SAMPLE_SID / "raw_response.txt"
            if p.exists():
                txt = p.read_text(errors="replace")
                # strip fences; keep first 600 chars
                txt = re.sub(r"```[a-z]*\n?", "", txt).strip()
                sample_outputs[model] = txt[:600]
                break

    # ── global counts ────────────────────────────────────────────────────────
    n_models    = ps[ps["rate"] > 0].shape[0]   # models with ≥1 parse
    n_models_all= ps.shape[0]
    n_configs   = len(configs_present)
    n_scenarios = df["scenario_id"].nunique()
    n_runs      = len(df)
    n_metrics   = len(df.columns)

    # ── winner analysis ──────────────────────────────────────────────────────
    # best parse rate, best F1, best hallucination (lowest)
    winners = {}
    if not qual.empty:
        if "combined_f1" in qual.columns:
            winners["best_f1"]    = qual["combined_f1"].idxmax()
        if "cq_coverage_ratio" in qual.columns:
            winners["best_cq"]    = qual["cq_coverage_ratio"].idxmax()
        if "hallucination_ratio" in qual.columns:
            winners["best_halluc"]= qual["hallucination_ratio"].idxmin()
    winners["best_parse"] = ps["rate"].idxmax()

    return dict(
        df=df, success=success,
        ps=ps, qual=qual, cfg_qual=cfg_qual, cross=cross,
        models_ordered=models_ordered,
        configs_present=configs_present,
        scenarios=scenarios,
        cutoffs=cutoffs,
        sample_outputs=sample_outputs,
        n_models=n_models, n_models_all=n_models_all,
        n_configs=n_configs, n_scenarios=n_scenarios,
        n_runs=n_runs, n_metrics=n_metrics,
        avail_q=avail_q,
        winners=winners,
        func_rank=func_rank,
    )


def _m(model_id: str, key: str, fallback="") -> str:
    """Look up MODEL_META; auto-generate fallback for unknown models."""
    meta = MODEL_META.get(model_id)
    if meta:
        return meta.get(key, fallback)
    # auto-generate display name from model_id
    if key == "name":
        return model_id.split("/")[-1].replace("_", " ")
    if key == "chip":
        return "chip-blue"
    if key == "pptx_color":
        return (0x8a, 0xb4, 0xf8)
    return fallback


def _fmt(val, decimals=3):
    """Format float or return '—'."""
    try:
        f = float(val)
        return f"{f:.{decimals}f}"
    except Exception:
        return "—"


def _pct(val):
    try:
        return f"{float(val)*100:.1f}%"
    except Exception:
        return "—%"


# ═══════════════════════════════════════════════════════════════════════════════
# HTML GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

def _color_val(val, col, q):
    """Wrap a metric value in a <span class=good/mid/bad> based on column semantics."""
    try:
        f = float(val)
    except Exception:
        return val
    # Higher-is-better columns
    if col in ("parse_rate","cq_coverage_ratio","scenario_coverage_ratio",
               "class_label_ratio","property_label_ratio","combined_f1",
               "class_name_f1","property_name_f1"):
        if q is None:
            return val
        hi = q.get(col, {}).get("max", 1)
        lo = q.get(col, {}).get("min", 0)
        rng = hi - lo if hi != lo else 1
        norm = (f - lo) / rng
        cls = "good" if norm >= 0.6 else ("bad" if norm <= 0.3 else "mid")
        return f'<span class="{cls}">{val}</span>'
    # Lower-is-better
    if col in ("hallucination_ratio",):
        hi = q.get(col, {}).get("max", 1) if q else 1
        lo = q.get(col, {}).get("min", 0) if q else 0
        rng = hi - lo if hi != lo else 1
        norm = (f - lo) / rng
        cls = "good" if norm <= 0.3 else ("bad" if norm >= 0.7 else "mid")
        return f'<span class="{cls}">{val}</span>'
    return val


CSS = r"""
  :root {
    --bg:#0f1117;--surface:#1a1d2e;--surface2:#232740;
    --accent:#6c8ebf;--accent2:#8ab4f8;
    --green:#4caf7d;--yellow:#f0c060;--red:#e06c6c;--purple:#b39ddb;
    --text:#e8eaf0;--muted:#8892a4;--border:#2e3350;
    --font:'Segoe UI',system-ui,sans-serif;
    --mono:'Cascadia Code','Fira Code',monospace;
  }
  *{box-sizing:border-box;margin:0;padding:0;}
  body{font-family:var(--font);background:var(--bg);color:var(--text);line-height:1.6;}
  nav{position:fixed;top:0;left:0;right:0;z-index:100;
      background:rgba(15,17,23,.95);border-bottom:1px solid var(--border);
      padding:10px 24px;display:flex;align-items:center;gap:12px;
      backdrop-filter:blur(8px);flex-wrap:wrap;}
  nav span{font-size:12px;color:var(--muted);margin-right:4px;}
  nav a{font-size:12px;color:var(--accent2);text-decoration:none;
        padding:4px 10px;border-radius:4px;transition:background .2s;}
  nav a:hover{background:var(--surface2);}
  .slide{min-height:100vh;padding:100px 60px 60px;border-bottom:2px solid var(--border);
         max-width:1200px;margin:0 auto;}
  #title{display:flex;flex-direction:column;justify-content:center;align-items:flex-start;
         min-height:100vh;background:radial-gradient(ellipse at 30% 50%,#1a2240 0%,var(--bg) 70%);}
  #title .badge{font-size:12px;text-transform:uppercase;letter-spacing:2px;
                color:var(--accent);border:1px solid var(--accent);
                padding:4px 12px;border-radius:20px;margin-bottom:24px;}
  #title h1{font-size:64px;font-weight:700;line-height:1.1;
            background:linear-gradient(135deg,#e8eaf0 40%,var(--accent2));
            -webkit-background-clip:text;-webkit-text-fill-color:transparent;margin-bottom:16px;}
  #title .subtitle{font-size:22px;color:var(--muted);margin-bottom:48px;}
  #title .meta{display:flex;gap:32px;flex-wrap:wrap;}
  #title .meta-label{font-size:11px;text-transform:uppercase;letter-spacing:1px;color:var(--muted);}
  #title .meta-value{font-size:16px;color:var(--accent2);}
  .slide-header{display:flex;align-items:center;gap:16px;margin-bottom:40px;
                padding-bottom:16px;border-bottom:1px solid var(--border);}
  .slide-num{font-size:12px;color:var(--muted);background:var(--surface2);
             padding:4px 10px;border-radius:4px;font-family:var(--mono);}
  h2{font-size:36px;color:var(--accent2);}
  h3{font-size:22px;color:var(--accent2);margin-bottom:12px;}
  h4{font-size:16px;color:var(--yellow);margin-bottom:8px;}
  .grid2{display:grid;grid-template-columns:1fr 1fr;gap:24px;}
  .grid3{display:grid;grid-template-columns:1fr 1fr 1fr;gap:20px;}
  .grid4{display:grid;grid-template-columns:repeat(4,1fr);gap:16px;}
  .card{background:var(--surface);border:1px solid var(--border);border-radius:12px;padding:20px;}
  .card.accent{border-color:var(--accent);}
  .card.green{border-color:var(--green);}
  .card.yellow{border-color:var(--yellow);}
  table{width:100%;border-collapse:collapse;font-size:14px;}
  th{background:var(--surface2);padding:10px 14px;text-align:left;color:var(--accent2);
     border-bottom:2px solid var(--border);font-size:12px;text-transform:uppercase;letter-spacing:.5px;}
  td{padding:9px 14px;border-bottom:1px solid var(--border);color:var(--text);}
  tr:last-child td{border-bottom:none;}
  tr:hover td{background:var(--surface2);}
  .good{color:var(--green);font-weight:600;}
  .mid{color:var(--yellow);font-weight:600;}
  .bad{color:var(--red);font-weight:600;}
  .best{background:rgba(76,175,125,.12);}
  .stat-box{text-align:center;padding:24px 16px;background:var(--surface);
            border:1px solid var(--border);border-radius:12px;}
  .stat-num{font-size:48px;font-weight:700;color:var(--accent2);line-height:1;}
  .stat-label{font-size:13px;color:var(--muted);margin-top:8px;}
  .bar-chart{display:flex;flex-direction:column;gap:12px;}
  .bar-row{display:flex;align-items:center;gap:12px;}
  .bar-label{width:240px;font-size:13px;text-align:right;flex-shrink:0;}
  .bar-track{flex:1;background:var(--surface2);border-radius:4px;height:28px;overflow:hidden;}
  .bar-fill{height:100%;border-radius:4px;display:flex;align-items:center;padding-left:10px;
            font-size:12px;font-weight:700;color:#fff;}
  .bf-green{background:linear-gradient(90deg,var(--green),#66bb9a);}
  .bf-yellow{background:linear-gradient(90deg,var(--yellow),#f5d080);color:#222;}
  .bf-red{background:linear-gradient(90deg,var(--red),#e89090);}
  .bf-blue{background:linear-gradient(90deg,var(--accent),#8ab4f8);}
  .bf-purple{background:linear-gradient(90deg,var(--purple),#ce93d8);}
  pre{background:var(--surface);border:1px solid var(--border);border-radius:8px;
      padding:16px;font-family:var(--mono);font-size:12px;overflow-x:auto;
      color:#c3c8d8;line-height:1.6;max-height:220px;overflow-y:auto;}
  .callout{padding:14px 20px;border-radius:8px;border-left:4px solid var(--accent);
           background:rgba(108,142,191,.1);margin:16px 0;font-size:14px;}
  .callout.green{border-color:var(--green);background:rgba(76,175,125,.1);}
  .callout.yellow{border-color:var(--yellow);background:rgba(240,192,96,.1);}
  .callout.red{border-color:var(--red);background:rgba(224,108,108,.1);}
  .chip{display:inline-block;font-size:11px;padding:3px 10px;border-radius:12px;margin:2px;font-weight:600;}
  .chip-blue{background:rgba(108,142,191,.25);color:var(--accent2);}
  .chip-green{background:rgba(76,175,125,.25);color:var(--green);}
  .chip-yellow{background:rgba(240,192,96,.25);color:var(--yellow);}
  .chip-red{background:rgba(224,108,108,.25);color:var(--red);}
  .chip-purple{background:rgba(179,157,219,.25);color:var(--purple);}
  p{margin-bottom:10px;font-size:15px;}
  li{margin-bottom:6px;font-size:14px;}
  ul{padding-left:20px;}
  .muted{color:var(--muted);}
  .small{font-size:12px;}
  .mt16{margin-top:16px;}.mt24{margin-top:24px;}.mt32{margin-top:32px;}
  hr{border:none;border-top:1px solid var(--border);margin:24px 0;}
  @media print{.slide{page-break-after:always;min-height:unset;}nav{display:none;}}
"""


def build_html(s: dict) -> None:
    ps, qual, cfg_qual = s["ps"], s["qual"], s["cfg_qual"]
    models = s["models_ordered"]
    configs = s["configs_present"]
    avail_q = s["avail_q"]

    # ── range dict for colour-coding ────────────────────────────────────────
    q_range: dict = {}
    if not qual.empty:
        for c in avail_q:
            if c in qual.columns:
                q_range[c] = {"min": qual[c].min(), "max": qual[c].max()}

    # helper
    def chip(model_id, label=None):
        n  = label or _m(model_id, "name", model_id)
        cl = _m(model_id, "chip", "chip-blue")
        return f'<span class="chip {cl}">{n}</span>'

    def stat_box(num, label):
        return (f'<div class="stat-box"><div class="stat-num">{num}</div>'
                f'<div class="stat-label">{label}</div></div>')

    def bar(label, pct, text, cls="bf-green"):
        w = max(1, int(pct * 100))
        return (f'<div class="bar-row"><span class="bar-label">{label}</span>'
                f'<div class="bar-track"><div class="bar-fill {cls}" style="width:{w}%">{text}</div></div></div>')

    def callout(text, kind=""):
        return f'<div class="callout {kind}">{text}</div>'

    # ── build model parse bars ───────────────────────────────────────────────
    def bar_class(rate):
        if rate >= 0.7: return "bf-green"
        if rate >= 0.35: return "bf-yellow"
        return "bf-red"

    parse_bars = "\n".join(
        bar(_m(m,"name",m),
            ps.loc[m,"rate"],
            f"{_pct(ps.loc[m,'rate'])}  ({int(ps.loc[m,'parsed'])}/{int(ps.loc[m,'total'])})",
            bar_class(ps.loc[m,"rate"]))
        for m in models if m in ps.index
    )

    # ── config parse bars ────────────────────────────────────────────────────
    cfg_ps = (s["df"].groupby("config")["parse_success"]
                .agg(["sum","count","mean"]).rename(columns={"mean":"rate"}))
    config_bars = "\n".join(
        bar(c,
            cfg_ps.loc[c,"rate"],
            f"{_pct(cfg_ps.loc[c,'rate'])}  ({int(cfg_ps.loc[c,'sum'])}/{int(cfg_ps.loc[c,'count'])})",
            "bf-blue")
        for c in configs if c in cfg_ps.index
    )

    # ── model quality table rows ─────────────────────────────────────────────
    show_cols = ["cq_coverage_ratio","scenario_coverage_ratio","hallucination_ratio",
                 "class_label_ratio","property_label_ratio","combined_f1"]
    show_cols = [c for c in show_cols if c in avail_q]

    def q_headers():
        labels = {"cq_coverage_ratio":"CQ Coverage","scenario_coverage_ratio":"Scenario Cov.",
                  "hallucination_ratio":"Hallucination","class_label_ratio":"Class Label ★",
                  "property_label_ratio":"Prop Label ★","combined_f1":"Combined F1"}
        return "".join(f"<th>{labels.get(c,c)}</th>" for c in show_cols)

    def q_row(model_id):
        pr = _pct(ps.loc[model_id,"rate"]) if model_id in ps.index else "—"
        pr_cls = "good" if ps.loc[model_id,"rate"]>=0.7 else ("bad" if ps.loc[model_id,"rate"]==0 else "mid")
        cells = f'<td>{chip(model_id)}</td><td><span class="{pr_cls}">{pr}</span></td>'
        if not qual.empty and model_id in qual.index:
            for c in show_cols:
                val = _fmt(qual.loc[model_id, c]) if c in qual.columns else "—"
                cells += f"<td>{_color_val(val, c, q_range)}</td>"
        else:
            cells += "".join("<td>—</td>" for _ in show_cols)
        return f"<tr>{'<td class=best>' if model_id==s['winners'].get('best_parse') else '<td>'}{cells[4:]}" \
               if False else f"<tr>{cells}</tr>"

    model_quality_rows = "\n".join(q_row(m) for m in models)

    # ── config quality table ─────────────────────────────────────────────────
    def cfg_row(cfg):
        cells = f"<td><strong>{cfg}</strong></td>"
        if not cfg_qual.empty and cfg in cfg_qual.index:
            for c in show_cols:
                val = _fmt(cfg_qual.loc[cfg, c]) if c in cfg_qual.columns else "—"
                cells += f"<td>{_color_val(val, c, q_range)}</td>"
        else:
            cells += "".join("<td>—</td>" for _ in show_cols)
        return f"<tr>{cells}</tr>"

    config_quality_rows = "\n".join(cfg_row(c) for c in configs)

    # ── heatmap table ────────────────────────────────────────────────────────
    cross = s["cross"]

    def heat_cell(model, cfg):
        if (model, cfg) not in cross.index:
            return '<td class="muted small">—</td>'
        r = cross.loc[(model, cfg), "rate"]
        cls = "good" if r >= 0.7 else ("bad" if r == 0 else "mid")
        return f'<td><span class="{cls}">{_pct(r)}</span></td>'

    heat_rows = ""
    for m in models:
        cells = "".join(heat_cell(m, c) for c in configs)
        avg = ps.loc[m,"rate"] if m in ps.index else 0
        avg_cls = "good" if avg>=0.7 else("bad" if avg==0 else "mid")
        heat_rows += (f"<tr><td>{chip(m)}</td>{cells}"
                      f'<td><strong><span class="{avg_cls}">{_pct(avg)}</span></strong></td></tr>\n')

    # ── sample outputs ───────────────────────────────────────────────────────
    sample_cards = ""
    for m in models:
        txt = s["sample_outputs"].get(m, "(no output found)")
        escaped = txt.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
        name = _m(m,"name",m)
        ok_flag = "SUCCESS" if m in ps.index and ps.loc[m,"rate"]>0 else "FAIL"
        ok_cls = "good" if ok_flag=="SUCCESS" else "bad"
        sample_cards += f"""
<div class="card" style="margin-bottom:16px">
  <h4>{chip(m)} — scenario-cq — <span class="{ok_cls}">{ok_flag}</span></h4>
  <pre>{escaped[:500]}</pre>
</div>"""

    # ── winner analysis ──────────────────────────────────────────────────────
    wp = _m(s["winners"].get("best_parse",""), "name", s["winners"].get("best_parse",""))
    wf = _m(s["winners"].get("best_f1",""),    "name", s["winners"].get("best_f1",""))
    wh = _m(s["winners"].get("best_halluc",""),"name", s["winners"].get("best_halluc",""))

    # ── scenarios table ──────────────────────────────────────────────────────
    scenario_rows = ""
    for sc in s["scenarios"]:
        yr = sc["year"]
        sid = sc["id"]
        domain = sc.get("domain", "—")
        cqs = sc["cqs"]
        if yr <= 2023:
            tag = '<span class="chip chip-green">Likely Seen</span>'
        elif yr == 2024:
            tag = '<span class="chip chip-yellow">Borderline</span>'
        else:
            tag = '<span class="chip chip-red">Plausible Unseen</span>'
        scenario_rows += f"<tr><td><strong>{sid}</strong></td><td>{domain}</td><td style='text-align:center'>{cqs}</td><td>{yr}</td><td>{tag}</td></tr>\n"

    # ── functional eval table rows ───────────────────────────────────────────
    _fs_max = max((r["avg_functional_score"] for r in s["func_rank"]), default=1)
    _func_table_rows = ""
    for _fr in s["func_rank"]:
        _fs = _fr["avg_functional_score"]
        _is_best = _fs >= _fs_max - 0.001
        _row_cls = ' class="best"' if _is_best else ""
        _fs_cls = "good" if _is_best else ("mid" if _fs >= 0.45 else "bad")
        _ms_cls = "good" if _fr["avg_manual_score"] >= 0.82 else ("mid" if _fr["avg_manual_score"] >= 0.70 else "bad")
        _cq_cls = "good" if _fr["avg_cq_pass_rate"] >= 0.18 else ("mid" if _fr["avg_cq_pass_rate"] >= 0.10 else "bad")
        _func_table_rows += (
            f'<tr{_row_cls}><td><strong>{_fr["model"]}</strong></td>'
            f'<td class="muted">{_fr["config"]}</td>'
            f'<td class="muted">{_fr["avg_numb"]:.2f}</td>'
            f'<td><span class="{_ms_cls}">{_fr["avg_manual_score"]:.3f}</span></td>'
            f'<td><span class="{_cq_cls}">{_fr["avg_cq_pass_rate"]:.3f}</span></td>'
            f'<td><strong><span class="{_fs_cls}">{_fr["avg_functional_score"]:.3f}</span></strong></td>'
            f'</tr>\n'
        )

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width,initial-scale=1.0">
<title>ODPGen — LLM-Based Ontology Design Pattern Generation</title>
<style>{CSS}</style>
</head>
<body>
<nav>
  <span>ODPGen</span>
  <a href="#title">Title</a>
  <a href="#overview">Overview</a>
  <a href="#dataset">Dataset</a>
  <a href="#example">Worked Example</a>
  <a href="#prompttext">Prompt Text</a>
  <a href="#models">Models</a>
  <a href="#prompts">Prompts</a>
  <a href="#evaluation">Evaluation</a>
  <a href="#outputs">Outputs</a>
  <a href="#results">Results</a>
  <a href="#quality">Quality</a>
  <a href="#heatmap">Heatmap</a>
  <a href="#structural">Structural Rankings</a>
  <a href="#functional">Functional Eval</a>
  <a href="#openmodel">Open Models</a>
  <a href="#humaneval">Human Eval Plan</a>
  <a href="#matched">Matched Names</a>
  <a href="#summary">Summary</a>
</nav>

<!-- SLIDE 1 — TITLE -->
<div class="slide" id="title">
  <div class="badge">Research Presentation</div>
  <h1>ODPGen</h1>
  <p class="subtitle">Benchmarking Large Language Models on<br>Ontology Design Pattern Generation</p>
  <div class="meta">
    <div><div class="meta-label">Models Evaluated</div><div class="meta-value">{s['n_models_all']} LLMs</div></div>
    <div><div class="meta-label">Prompt Strategies</div><div class="meta-value">{s['n_configs']} Configurations</div></div>
    <div><div class="meta-label">ODP Scenarios</div><div class="meta-value">{s['n_scenarios']} Benchmarks</div></div>
    <div><div class="meta-label">Total Runs</div><div class="meta-value">{s['n_runs']} Evaluated · {s['n_metrics']} Metrics</div></div>
  </div>
</div>

<!-- SLIDE 2 — OVERVIEW -->
<div class="slide" id="overview">
  <div class="slide-header"><span class="slide-num">01</span><h2>Overview &amp; Research Questions</h2></div>
  <p>ODPGen evaluates how well LLMs generate <strong>Ontology Design Patterns (ODPs)</strong> — reusable OWL ontology templates — from natural-language scenarios and competency questions.</p>
  <div class="grid2 mt24">
    <div class="card accent">
      <h3>Research Questions</h3>
      <ul>
        <li><strong>RQ1</strong> — How well can LLMs generate valid OWL/Turtle from scenarios?</li>
        <li><strong>RQ2</strong> — How logically consistent and error-free are the outputs?</li>
        <li><strong>RQ3</strong> — What is the impact of prompting strategy?</li>
        <li><strong>RQ4</strong> — How faithful are outputs to the ground-truth scenario?</li>
        <li><strong>RQ5</strong> — Can models generalize to recently published ODPs?</li>
      </ul>
    </div>
    <div>
      <div class="grid2">
        {stat_box(s['n_models_all'], 'LLMs Evaluated')}
        {stat_box(s['n_configs'], 'Prompt Configs')}
        {stat_box(s['n_scenarios'], 'ODP Scenarios')}
        {stat_box(s['n_runs'], 'Total Runs')}
      </div>
    </div>
  </div>
  {callout('<strong>Core hypothesis:</strong> Structured prompting (CQs + reasoning or self-validation) improves faithfulness — but may trade off syntactic correctness for weaker models.')}
</div>

<!-- SLIDE 3 — DATASET -->
<div class="slide" id="dataset">
  <div class="slide-header"><span class="slide-num">02</span><h2>Dataset — {s['n_scenarios']} ODP Scenarios</h2></div>
  <p>Scenarios sourced from published ODP papers (2023–2026), each with a natural-language domain description, Competency Questions (CQs), and a ground-truth OWL ontology. CQ counts range from 3 to 47 per scenario, reflecting varying pattern complexity.</p>
  <table class="mt24">
    <thead><tr><th>Scenario ID</th><th>Domain</th><th># CQs</th><th>Year</th><th>Unseen Status</th></tr></thead>
    <tbody>{scenario_rows}</tbody>
  </table>
  <div class="grid2 mt24">
    {callout('<strong>Unseen analysis:</strong> 9 of 14 scenarios were published after most model training cutoffs — enabling memorization vs. generalization analysis. Scenarios from 2025–2026 are <em>plausible unseen</em> for all evaluated models.', 'yellow')}
    {callout('<strong>CQ range:</strong> 3 CQs (Role-Dependent Names) to 47 CQs (AISHIP Maritime). Very long CQ lists may be trimmed for human evaluation to avoid participant fatigue.', 'green')}
  </div>
</div>

<!-- SLIDE 3b — WORKED EXAMPLE -->
<div class="slide" id="example">
  <div class="slide-header"><span class="slide-num">03</span><h2>Worked Example — Scenario 2023-133-01: Causal Bayesian Network</h2></div>
  <div class="grid2">
    <div>
      <div class="card accent" style="margin-bottom:16px">
        <h4>Scenario Text</h4>
        <p class="muted small">The sprinkler use case is a well-known example from the CBN literature. A CBN reasons about the cause of a slippery pavement. Five nodes represent events: <em>rainy season</em>, <em>sprinkler on</em>, <em>raining</em>, <em>pavement wet</em>, <em>pavement slippery</em>. Directed edges carry effect weights — e.g., rainy season → rain with weight 0.6. Rain→slippery is mediated by pavement being wet.</p>
      </div>
      <div class="card" style="margin-bottom:16px">
        <h4>Competency Questions (4 CQs)</h4>
        <ol style="padding-left:18px;font-size:13px;color:var(--text)">
          <li>What events are responsible for the pavement being slippery?</li>
          <li>What would be the outcome of turning the sprinkler on/off?</li>
          <li>How strong is the causal link between sprinkler and wet pavement?</li>
          <li>If the pavement is not wet, how does that affect it being slippery?</li>
        </ol>
      </div>
      <div class="card green">
        <h4>Ground-Truth ODP Structure (Expert-Authored)</h4>
        <p class="muted small" style="margin-bottom:8px">5 classes · 5 object properties · reification pattern</p>
        <ul style="font-size:12px;font-family:var(--mono)">
          <li>Classes: <span class="good">Event</span>, AbstractEvent, ConcreteEvent, <span class="good">Causes</span>, EffectWeight</li>
          <li>hasTreatment, hasOutcome (Causes→Event)</li>
          <li>hasMediator (Causes→Event — captures indirect causation)</li>
          <li>hasEffectWeight (Causes→EffectWeight)</li>
          <li>ofType (ConcreteEvent→AbstractEvent)</li>
        </ul>
        <p class="muted small mt16">Key design: <strong>Causes is a class</strong> (reification), not a property — enables weight and mediator to be attached to the relation itself.</p>
      </div>
    </div>
    <div>
      <h3>Generated Outputs — scenario-only</h3>
      <div class="card" style="margin-bottom:12px;border-color:var(--accent2)">
        <h4><span class="chip chip-blue">Gemini 3.1 Pro</span> — Parse: <span class="good">SUCCESS</span></h4>
        <pre style="font-size:11px;max-height:160px">:Event a owl:Class .
:CausalRelation a owl:Class ; owl:disjointWith :Event .
:RainySeasonEvent rdfs:subClassOf :Event .
:SprinklerOnEvent rdfs:subClassOf :Event .
:PavementWetEvent rdfs:subClassOf :Event .
:PavementSlipperyEvent rdfs:subClassOf :Event .
:hasSourceEvent, :hasTargetEvent, :isMediatedBy
    rdfs:domain :CausalRelation ; rdfs:range :Event .
:hasEffectWeight a owl:DatatypeProperty ;
    rdfs:range xsd:decimal .</pre>
        <p class="muted small">+ 5 domain-specific event subclasses (scenario-grounded but not in ground truth). Does not use reification pattern — models causal weight as datatype property of relation class.</p>
      </div>
      <div class="card" style="margin-bottom:12px;border-color:var(--green)">
        <h4><span class="chip chip-green">Llama-3.1-8B</span> — Parse: <span class="good">SUCCESS</span></h4>
        <pre style="font-size:11px;max-height:130px">:WeatherEvent, :RainySeason, :SprinklerStatus,
:SprinklerOn, :Weather, :Rain,
:PavementCondition, :PavementWet, :PavementSlippery
:CausalRelation a owl:Class .
:EffectWeight a owl:DatatypeProperty .</pre>
        <p class="muted small">Different conceptualization: introduces intermediate categories (WeatherEvent, PavementCondition) not in ground truth or scenario. More superfluous elements.</p>
      </div>
      <div class="callout yellow">
        <strong>Systematic finding:</strong> Both models literalize the sprinkler scenario into domain-specific class hierarchies. The expert-authored ODP is <em>abstract and reusable</em> (no sprinkler/rain classes) — it uses reification to make causal relations first-class objects. LLMs struggle to generate this level of abstraction from narrative alone.
      </div>
    </div>
  </div>
</div>

<!-- SLIDE 3c — PROMPT TEMPLATES FULL TEXT -->
<div class="slide" id="prompttext">
  <div class="slide-header"><span class="slide-num">04</span><h2>Prompt Templates — Actual Text Sent to Models</h2></div>
  <div class="grid2">
    <div>
      <div class="card accent">
        <h4>Config 1 — scenario-only</h4>
        <pre style="font-size:10px;max-height:340px">You are an ontology engineer. Generate an OWL
ontology design pattern (ODP) in Turtle syntax
from the scenario below.

Requirements:
- Stay strictly within the scenario scope.
- Use clear, self-explanatory class/property names.
- Include minimal but sufficient axioms.
- Declare every class with `a owl:Class`, every
  property with the appropriate OWL type.
- Do NOT create individual instances.
- Add rdfs:label and rdfs:comment to all entities.
- Declare: :Ontology a owl:Ontology ; rdfs:label
  "..." ; rdfs:comment "..." .

@prefix owl:  &lt;...owl#&gt; .
@prefix rdf:  &lt;...rdf-syntax-ns#&gt; .
@prefix rdfs: &lt;...rdf-schema#&gt; .
@prefix xsd:  &lt;...XMLSchema#&gt; .
@prefix :     &lt;http://example.org/odp#&gt; .

Output — in this order:
1. Turtle ontology in ```turtle ... ``` block.
2. Documentation: Intent · Requirements · Assumptions.

Scenario:
{{SCENARIO_TEXT}}</pre>
      </div>
    </div>
    <div>
      <div class="card green">
        <h4>Config 5 — scenario-cq-reasoning (best faithfulness)</h4>
        <pre style="font-size:10px;max-height:340px">You are an ontology engineer. Create an OWL ODP
in Turtle syntax.

Before writing, reason through these steps
internally (do NOT write them out yet):
1. Extract atomic requirements from the scenario.
2. Propose candidate classes/properties per req.
3. Draft axioms (subclass, domain, range, card.).
4. For each CQ, verify which axioms answer it.
5. Add/refine axioms to cover uncovered CQs.
6. Remove entities not traceable to any req/CQ.

Constraints:
- schema axioms only (no instances)
- rdfs:label + rdfs:comment on every entity
- owl:Ontology header mandatory
- minimal and reusable; no facts beyond scenario/CQs

Output — strictly in this order:
1. Complete Turtle in ```turtle ... ``` block.
2. Reasoning summary: one bullet per step 1–6.
3. Traceability table: requirement/CQ → axiom(s).

Scenario: {{SCENARIO_TEXT}}
CQs: {{CQ_LIST}}</pre>
      </div>
      <div class="callout mt16">
        <strong>Config 3 (scenario-cq)</strong> adds scenario→axiom and CQ→axiom mapping tables to the output. <strong>Config 4 (scenario-cq-constraints)</strong> adds a 6-item self-validation checklist (syntax, prefixes, domain/range, coherence, CQ coverage, scope). All configs enforce temperature=0.0.
      </div>
    </div>
  </div>
</div>

<!-- SLIDE 4 — MODELS -->
<div class="slide" id="models">
  <div class="slide-header"><span class="slide-num">03</span><h2>Models Evaluated</h2></div>
  <table>
    <thead><tr><th>Model</th><th>Provider</th><th>Parameters</th><th>Training Cutoff</th><th>Parse Rate</th></tr></thead>
    <tbody>
{"".join(f'''<tr>
  <td>{chip(m)}</td>
  <td class="muted">{_m(m,"provider","—")}</td>
  <td class="muted">{_m(m,"params","—")}</td>
  <td class="muted">{_m(m,"cutoff","—")}</td>
  <td><span class="{"good" if ps.loc[m,"rate"]>=0.7 else ("bad" if ps.loc[m,"rate"]==0 else "mid")}">{_pct(ps.loc[m,"rate"])}</span></td>
</tr>''' for m in models if m in ps.index)}
    </tbody>
  </table>
  {callout('All runs use temperature=0.0. Open-source models run via Hugging Face Transformers; GPT-5.4 via OpenAI API; Gemini via Google GenAI SDK.')}
</div>

<!-- SLIDE 5 — PROMPTS -->
<div class="slide" id="prompts">
  <div class="slide-header"><span class="slide-num">04</span><h2>Prompting Strategies — {s['n_configs']} Configurations</h2></div>
  <div class="grid2" style="gap:20px">
    <div class="card"><h4>1. scenario-only</h4><p>Model receives only the domain scenario text. No CQs. Must infer scope from narrative alone. Highest scenario vocabulary coverage.</p></div>
    <div class="card"><h4>2. cq-only</h4><p>Model receives only the Competency Questions. No scenario text. Highest CQ coverage but also highest hallucination.</p></div>
    <div class="card green"><h4>3. scenario-cq &nbsp;<small>[Balanced]</small></h4><p>Both scenario text and CQs. Model must produce scenario→axiom and CQ→axiom mapping tables.</p></div>
    <div class="card yellow"><h4>4. scenario-cq-constraints &nbsp;<small>[Validation]</small></h4><p>Extends scenario-cq with a 6-item self-validation checklist before output.</p></div>
    <div class="card" style="grid-column:span 2"><h4>5. scenario-cq-reasoning &nbsp;<small>[Best Faithfulness]</small></h4><p>Explicit 6-step internal reasoning: extract requirements → propose classes → draft axioms → verify CQs → remove unsupported → output ODP. Lowest hallucination overall. Also produces a reasoning summary and traceability table (requirement/CQ → axiom).</p></div>
  </div>
</div>

<!-- SLIDE 6 — EVALUATION -->
<div class="slide" id="evaluation">
  <div class="slide-header"><span class="slide-num">05</span><h2>Evaluation Methodology — 4 Layers</h2></div>
  <div class="grid2">
    <div class="card accent">
      <h4>Layer 1 — Formal / Syntactic Quality</h4>
      <ul>
        <li><strong>parse_success</strong> — Valid Turtle syntax?</li>
        <li><strong>triple_count</strong> — Total RDF triples</li>
        <li><strong>ontology_declared</strong> — owl:Ontology header?</li>
        <li><strong>label_coverage_ratio</strong> — All entities with rdfs:label</li>
        <li><strong>class_label_ratio</strong>  — owl:Class entities with rdfs:label</li>
        <li><strong>property_label_ratio</strong>  — Properties with rdfs:label</li>
        <li><strong>comment_coverage_ratio, domain_range_ratio</strong></li>
      </ul>
      {callout('<strong>Why split?</strong> The aggregate label_coverage_ratio was always 1.0 (dominated by classes). Splitting exposes that property labelling is inconsistent — Llama-2-70B leaves 28% of properties unlabelled.')}
    </div>
    <div class="card green">
      <h4>Layer 2 — Scenario Faithfulness</h4>
      <ul>
        <li><strong>cq_coverage_ratio</strong> — CQs answerable (token overlap)</li>
        <li><strong>cq_mean_token_recall</strong> — Avg token recall per CQ</li>
        <li><strong>scenario_coverage_ratio</strong> — Scenario vocab in ontology</li>
        <li><strong>hallucination_ratio</strong> — Unsupported vocabulary fraction</li>
      </ul>
    </div>
    <div class="card yellow">
      <h4>Layer 3 — Semantic Similarity (vs Ground Truth)</h4>
      <ul>
        <li><strong>class_name_F1</strong> — Name match vs ground-truth classes</li>
        <li><strong>property_name_F1</strong> — Same for properties</li>
        <li><strong>combined_f1</strong> — Macro avg F1</li>
        <li><strong>matched_class_names</strong> — Classes in both generated &amp; GT</li>
        <li><strong>matched_property_names</strong> — Properties in both</li>
      </ul>
      {callout('Matching: local name extracted, CamelCase split, case-folded. :hasCause ↔ "has cause".')}
    </div>
    <div class="card" style="border-color:var(--purple)">
      <h4>Layer 4 — Feasibility &amp; Unseen Analysis</h4>
      <ul>
        <li>Each (ODP, model) pair classified: <em>plausible_unseen</em> vs <em>likely_seen</em></li>
        <li>Compares ODP pub date vs model training cutoff</li>
        <li>Confidence: high / medium / low</li>
        <li>Future: human eval — Correctness, Completeness, Clarity (1–5)</li>
      </ul>
    </div>
  </div>
  {callout(f'Pipeline: evaluate_outputs.py → eval/formal_quality.py + faithfulness.py + similarity.py → results/summary.csv  ({s["n_runs"]} rows × {s["n_metrics"]} columns)')}
</div>

<!-- SLIDE 7 — OUTPUTS -->
<div class="slide" id="outputs">
  <div class="slide-header"><span class="slide-num">06</span><h2>Raw Output Examples — Scenario 2023-133-01 (Causal Bayesian Network)</h2></div>
  <p class="muted small mb16">Scenario: "Model a causal Bayesian network over events with weighted causal edges and mediators as an ODP."</p>
  {sample_cards}
</div>

<!-- SLIDE 8 — PARSE RESULTS -->
<div class="slide" id="results">
  <div class="slide-header"><span class="slide-num">07</span><h2>Results — Parse Success Rate</h2></div>
  <div class="grid2">
    <div>
      <h3>By Model</h3>
      <div class="bar-chart mt16">{parse_bars}</div>
    </div>
    <div>
      <h3>By Configuration</h3>
      <div class="bar-chart mt16">{config_bars}</div>
    </div>
  </div>
  {callout('<strong>Gemini 3.1 Pro</strong> achieves <strong>100%</strong> parse success across all 5 configurations — the only model to do so. Llama-3.1-8B is the best open-source model at 75.7%.','green')}
  {callout('GPT-5.4 paradox: 86% on cq-only but 0% on reasoning/constraints configs — complex prompts break its Turtle extraction.','yellow')}
</div>

<!-- SLIDE 9 — QUALITY METRICS -->
<div class="slide" id="quality">
  <div class="slide-header"><span class="slide-num">08</span><h2>Results — Quality Metrics (Successful Parses)</h2></div>
  <p class="muted small mb16">Averages over successful parses only. {s['n_runs']} total runs · {s['n_metrics']} metrics.</p>
  <table>
    <thead>
      <tr><th>Model</th><th>Parse Rate</th>{q_headers()}</tr>
    </thead>
    <tbody>{model_quality_rows}</tbody>
  </table>
  {callout('<strong>class_label_ratio = 1.000</strong> for every model — classes always annotated. <strong>property_label_ratio</strong> reveals the gap: Llama-2-70B 0.717 (worst), Mistral 0.966 (best). Old aggregate label_coverage_ratio masked this.','yellow')}
  <hr>
  <h3 class="mt24">By Configuration</h3>
  <table class="mt16">
    <thead>
      <tr><th>Config</th>{q_headers()}</tr>
    </thead>
    <tbody>{config_quality_rows}</tbody>
  </table>
</div>

<!-- SLIDE 10 — HEATMAP -->
<div class="slide" id="heatmap">
  <div class="slide-header"><span class="slide-num">09</span><h2>Model × Config Parse Success Heatmap</h2></div>
  <table>
    <thead>
      <tr><th>Model</th>{"".join(f"<th>{c}</th>" for c in configs)}<th>Average</th></tr>
    </thead>
    <tbody>{heat_rows}</tbody>
  </table>
  {callout('Gemini 3.1 Pro: 100% across every single config — the only model with perfect reliability.','green')}
  {callout('BLOOMZ-7B1: 0% across all configs — instruction-tuned BLOOM still cannot produce valid Turtle.','red')}
</div>

<!-- SLIDE 9b — STRUCTURAL RANKINGS -->
<div class="slide" id="structural">
  <div class="slide-header"><span class="slide-num">10</span><h2>Pilot Results — Structural Evaluation Rankings</h2></div>
  <p>Structural score = 0.5 × parse_rate + 0.5 × avg_OOPS_score. Computed per (model, config) pair across 14 scenarios. OOPS score measures freedom from common OWL modeling pitfalls (0 = worst, 1 = best).</p>
  <div class="grid2 mt24">
    <div>
      <h3>Proprietary Models — Top 4 Configurations</h3>
      <table class="mt16">
        <thead><tr><th>Rank</th><th>Model</th><th>Config</th><th>Struct. Score</th><th>OOPS Score</th><th>Parse Rate</th></tr></thead>
        <tbody>
          <tr class="best"><td><strong>1</strong></td><td><span class="chip chip-blue">Gemini 3.1 Pro</span></td><td>scenario-only</td><td><span class="good">0.982</span></td><td><span class="good">0.964</span></td><td><span class="good">100%</span></td></tr>
          <tr><td>2</td><td><span class="chip chip-blue">Gemini 3.1 Pro</span></td><td>scenario-cq-reasoning</td><td><span class="good">0.857</span></td><td><span class="good">0.714</span></td><td><span class="good">100%</span></td></tr>
          <tr><td>3</td><td><span class="chip chip-purple">GPT-5.4</span></td><td>scenario-only</td><td><span class="mid">0.464</span></td><td><span class="mid">0.429</span></td><td><span class="mid">50%</span></td></tr>
          <tr><td>4</td><td><span class="chip chip-purple">GPT-5.4</span></td><td>scenario-cq-constraints</td><td><span class="bad">0.293</span></td><td><span class="bad">0.214</span></td><td><span class="bad">37%</span></td></tr>
        </tbody>
      </table>
      <div class="callout green mt16"><strong>Gemini scenario-only (0.982)</strong> is the top structural result across ALL models and configurations — outperforming even scenario-cq-reasoning (0.857). Adding CQs does not always improve structural quality.</div>
    </div>
    <div>
      <h3>Open-Source Models — Top 3 Configurations</h3>
      <table class="mt16">
        <thead><tr><th>Rank</th><th>Model</th><th>Config</th><th>Struct. Score</th><th>OOPS Score</th><th>Parse Rate</th></tr></thead>
        <tbody>
          <tr><td>1</td><td><span class="chip chip-green">Llama-3.1-8B</span></td><td>scenario-cq</td><td><span class="mid">0.616</span></td><td><span class="mid">0.375</span></td><td><span class="good">85.7%</span></td></tr>
          <tr class="best"><td><strong>2★</strong></td><td><span class="chip chip-green">Llama-3.1-8B</span></td><td>scenario-only</td><td><span class="mid">0.595</span></td><td><span class="mid">0.405</span></td><td><span class="mid">78.6%</span></td></tr>
          <tr><td>3★</td><td><span class="chip chip-yellow">Llama-2-70B</span></td><td>cq-only</td><td><span class="mid">0.589</span></td><td><span class="mid">0.321</span></td><td><span class="good">85.7%</span></td></tr>
        </tbody>
      </table>
      <div class="callout yellow mt16"><strong>★ Selected for human evaluation</strong> — Llama-3.1-8B scenario-only (structural/creativity track) and Llama-2-70B cq-only (functional track). Though scenario-cq ranks #1 structurally for Llama-8B, scenario-only is selected to test prompt minimalism.</div>
      <div class="callout red mt16"><strong>Anomaly:</strong> Llama-3.1-8B scenario-cq on scenario 2025-151-01 produced 22 repetitions of identical triple blocks (degenerate output). Parseable but flagged — retained in dataset as an edge case.</div>
    </div>
  </div>
</div>

<!-- SLIDE 10 — FUNCTIONAL EVALUATION -->
<div class="slide" id="functional">
  <div class="slide-header"><span class="slide-num">10</span><h2>Functional Evaluation — Utility-Oriented Pilot Assessment</h2></div>
  <p>Functional evaluation pilots a utility-oriented composite score, measuring whether generated ODPs answer CQs without adding unsupported ontological commitments. Applied to the best configuration per model family.</p>
  <div class="grid2 mt24">
    <div class="card accent">
      <h4>Scoring Methodology</h4>
      <p style="font-family:var(--mono);font-size:13px;margin:8px 0">functional_score = 0.5 × manual_score + 0.5 × cq_pass_rate</p>
      <hr style="margin:12px 0;border-color:var(--border)">
      <ul>
        <li><strong>numb</strong>: total elements added beyond ground truth</li>
        <li><strong>numbnoscen</strong>: extra elements with no scenario grounding</li>
        <li><strong>superfluous_penalty</strong> = numbnoscen × 1.0 + (numb − numbnoscen) × 0.3</li>
        <li><strong>manual_score</strong> = 1 / (1 + 0.5 × superfluous_penalty)</li>
        <li><strong>cq_pass_rate</strong>: fraction of CQs answerable by the generated ontology</li>
      </ul>
      <p class="muted small mt16">Unsupported additions penalized at 1.0; scenario-justified additions at 0.3 — reflecting differential harm.</p>
    </div>
    <div>
      <h3>Pilot Results — Top Configuration per Model Family</h3>
      <table class="mt16">
        <thead><tr><th>Model</th><th>Config</th><th>Avg Extra</th><th>Manual Score</th><th>CQ Pass Rate</th><th>Functional Score</th></tr></thead>
        <tbody>{_func_table_rows}</tbody>
      </table>
      <div class="callout yellow mt16"><strong>Key finding:</strong> Llama-2-70B (cq-only) achieves the best functional score (0.534) despite structurally poor ontologies — because it produces extremely minimal outputs (avg. 0.79 extra elements). Parsimony ≠ ontological quality.</div>
    </div>
  </div>
  <div class="callout red mt16"><strong>Evaluation gap identified:</strong> Automated functional metrics reward minimalism but cannot detect modeling deficiencies such as incorrect domain/range, missing restrictions, or semantically incoherent axioms. Expert human evaluation is essential to assess real-world utility.</div>
</div>

<!-- SLIDE 11 — OPEN-SOURCE MODEL SELECTION -->
<div class="slide" id="openmodel">
  <div class="slide-header"><span class="slide-num">11</span><h2>Open-Source Model Selection for Human Evaluation</h2></div>
  <p>Structural and functional evaluations identify <em>different</em> open-source winners. Two configurations are selected for human evaluation to represent both quality dimensions.</p>
  <div class="grid2 mt24">
    <div class="card green">
      <h4>Track A — Structural / Creativity Winner</h4>
      <p style="font-size:18px;font-weight:700;margin:8px 0">Llama-3.1-8B + scenario-only</p>
      <ul class="mt16">
        <li>Structural score 0.595 — best open-source scenario-only result</li>
        <li>OOPS score 0.405 — highest among open-source scenario-only runs</li>
        <li>Parse rate 78.6% across 14 scenarios</li>
        <li>Tests whether scenario text alone suffices for ODP generation</li>
        <li>Richer ontological outputs; higher creativity</li>
        <li>Trade-off: more superfluous additions than minimal configs</li>
      </ul>
    </div>
    <div class="card yellow">
      <h4>Track B — Functional Winner</h4>
      <p style="font-size:18px;font-weight:700;margin:8px 0">Llama-2-70B + cq-only</p>
      <ul class="mt16">
        <li>Functional score 0.534 — highest among all evaluated configurations</li>
        <li>CQ pass rate 0.202 — highest open-source cq-only result</li>
        <li>Parse rate 85.7% in cq-only configuration</li>
        <li>Average 0.79 extra elements: extremely minimal outputs</li>
        <li>Structural quality known to be poor — "incredibly bad modeling"</li>
        <li>Human experts will assess whether functional advantage survives scrutiny</li>
      </ul>
    </div>
  </div>
  <div class="callout green mt24"><strong>Why two tracks?</strong> Structural metrics reward OWL correctness and expressivity; functional metrics reward minimalism and CQ answerability. The fact that they select <em>different</em> models is itself a scientific finding, motivating human evaluation to determine which dimension predicts real-world usefulness.</div>
  <div class="callout mt16"><strong>Evaluation design:</strong> Reducing to 2 open-source configurations (from 5) makes human evaluation feasible. The scenario-only vs. cq-only contrast also probes the extremes of the prompting spectrum for open-source LLMs.</div>
</div>

<!-- SLIDE 11b — HUMAN EVALUATION PLAN -->
<div class="slide" id="humaneval">
  <div class="slide-header"><span class="slide-num">12</span><h2>Human Evaluation Plan</h2></div>
  <p>Expert and student evaluation of selected LLM-generated ODPs. Blind design: evaluators are not told whether an ODP was generated by an LLM or a human. Three evaluator groups receive tailored forms.</p>
  <div class="grid3 mt24">
    <div class="card accent">
      <h4>Pattern Authors</h4>
      <ul>
        <li><strong>Intent match</strong> (1–5): Does the ODP model what the original pattern was designed to model?</li>
        <li><strong>Scope appropriateness</strong> (1–5): breadth — too broad / too narrow / right</li>
        <li><strong>Reuse readiness</strong> (1–5): ready for reuse without major modification?</li>
        <li><strong>Adoption threshold</strong>: Would you start from this draft rather than scratch? (Yes/Maybe/No)</li>
        <li><strong>Time saved</strong> vs from-scratch (0% / 1–25% / 26–50% / 51–75% / 76–100%)</li>
        <li><strong>Missing core elements</strong>: free text — essential classes/properties absent from the draft</li>
      </ul>
    </div>
    <div class="card green">
      <h4>Ontology Experts</h4>
      <ul>
        <li><strong>Overall helpfulness</strong> (1–5): useful as a starting point for ontology design?</li>
        <li><strong>Reusability</strong> (1–5): reusable or adaptable in related contexts?</li>
        <li><strong>Clarity / Documentation</strong> (1–5): clear and complete accompanying documentation?</li>
        <li><strong>Adoption threshold</strong>: start from this draft or from scratch? (Yes/Maybe/No)</li>
        <li><strong>Time saved</strong> (0% → 76–100%)</li>
        <li><strong>Free text</strong>: observations, concerns, suggestions</li>
      </ul>
    </div>
    <div class="card yellow">
      <h4>MSc / PhD Students</h4>
      <ul>
        <li><strong>Time to understand</strong> (minutes): includes reading scenario + CQs</li>
        <li><strong>Documentation clarity</strong> (1–5)</li>
        <li><strong>Naming clarity</strong> (1–5): intuitive and self-explanatory class/property names?</li>
        <li><strong>Adoption</strong>: start from this draft rather than scratch? (Yes/Maybe/No)</li>
        <li><strong>Time saved</strong> (0% → 76–100%)</li>
      </ul>
    </div>
  </div>
  <div class="grid2 mt24">
    <div class="card">
      <h4>Participant Pool</h4>
      <ul>
        <li>Pattern authors: ~2 (incl. 1 also a student)</li>
        <li>Ontology experts: ~5 (some also pattern authors)</li>
        <li>MSc/PhD students: ~2</li>
        <li>ISWC workshop: ~8 (experts + students)</li>
        <li><strong>Target: ~38 expert evaluators + ~38 students</strong></li>
        <li>Each evaluator rates ~3 ontologies; authors rate ~2</li>
        <li>Platform: <a href="https://github.com/ebrahimnorouzi/odp-platform" style="color:var(--accent2)">odp-platform</a></li>
      </ul>
    </div>
    <div>
      <div class="callout yellow"><strong>Open questions:</strong> Very long CQ lists (e.g., 47 CQs for AISHIP) may deter participation — considering trimming to ~3 representative CQs per scenario. Whether to include ground-truth ODPs as a calibration item.</div>
      <div class="callout green mt16"><strong>Blind design:</strong> Evaluators are told that ontologies could be either LLM-generated or human-authored. Source is not disclosed, reducing expectation bias.</div>
      <div class="callout mt16"><strong>ODP evaluated in pilot:</strong> Cultural Heritage Survey (2026-155-01) — one pattern author present for this ODP attended the pilot session.</div>
    </div>
  </div>
</div>

<!-- SLIDE 12 — MATCHED NAMES -->
<div class="slide" id="matched">
  <div class="slide-header"><span class="slide-num">12</span><h2>Matched Class &amp; Property Names vs. Ground Truth</h2></div>
  <p>Names found in both generated ontology and ground truth (CamelCase split, case-folded). Low F1 does not mean zero conceptual overlap.</p>
  <div class="callout yellow mt16">Matching: local name from URI → split CamelCase → lowercase → exact string match vs ground truth vocabulary.</div>
  <table class="mt24">
    <thead><tr><th>Model / Config / Scenario</th><th>Matched Classes</th><th>Matched Properties</th></tr></thead>
    <tbody>
      <tr><td>GPT-5.4 / cq-only / 2023-133-01</td><td>event</td><td>has effect, has outcome</td></tr>
      <tr><td>GPT-5.4 / cq-only / 2023-135-01</td><td>name usage</td><td>has name usage</td></tr>
      <tr><td>GPT-5.4 / cq-only / 2025-149-01</td><td>relational database, table, column</td><td>—</td></tr>
      <tr><td>Gemini 3.1 Pro / scenario-cq / 2023-133-01</td><td>event</td><td>—</td></tr>
      <tr><td>Gemini 3.1 Pro / cq-only / 2025-151-01</td><td>evaluation request</td><td>requested action, issued by</td></tr>
    </tbody>
  </table>
  <hr>
  <h3 class="mt24">Property Label Ratio — {_m(s['winners'].get('best_halluc',''),'name','?')} has best hallucination control</h3>
  <div class="bar-chart mt16">
{"".join(bar(_m(m,"name",m), qual.loc[m,"property_label_ratio"] if (not qual.empty and m in qual.index and "property_label_ratio" in qual.columns) else 0, _fmt(qual.loc[m,"property_label_ratio"]) if (not qual.empty and m in qual.index and "property_label_ratio" in qual.columns) else "—", bar_class(qual.loc[m,"property_label_ratio"]) if (not qual.empty and m in qual.index and "property_label_ratio" in qual.columns) else "bf-red") for m in models if not qual.empty and m in qual.index and ps.loc[m,"rate"]>0)}
  </div>
</div>

<!-- SLIDE 13 — SUMMARY -->
<div class="slide" id="summary">
  <div class="slide-header"><span class="slide-num">13</span><h2>Summary &amp; Best Performing Model</h2></div>
  <div class="grid2">
    <div>
      <div class="card green" style="margin-bottom:16px">
        <h4>★ Best Overall Reliability: {wp}</h4>
        <ul>
          <li>Parse success: {_pct(ps.loc[s['winners']['best_parse'],'rate']) if s['winners'].get('best_parse') in ps.index else '—'} — highest across all configs</li>
          <li>100% parse on every single configuration</li>
          <li>Best CQ coverage and Combined F1 when successful</li>
          <li>Trade-off: highest hallucination ratio (adds unsupported domain knowledge)</li>
        </ul>
      </div>
      <div class="card yellow" style="margin-bottom:16px">
        <h4>🥈 Best Faithfulness: {wh}</h4>
        <ul>
          <li>Lowest hallucination: {_fmt(qual.loc[s['winners']['best_halluc'],'hallucination_ratio']) if s['winners'].get('best_halluc') in qual.index and 'hallucination_ratio' in qual.columns else '—'}</li>
          <li>Best open-source model for scenario-faithful generation</li>
          <li>Recommended config: scenario-cq-reasoning</li>
        </ul>
      </div>
      <div class="card" style="margin-bottom:16px">
        <h4>🥉 Best Ground-Truth Match: {wf}</h4>
        <ul>
          <li>Combined F1: {_fmt(qual.loc[s['winners']['best_f1'],'combined_f1']) if s['winners'].get('best_f1') in qual.index and 'combined_f1' in qual.columns else '—'}</li>
          <li>Closest to expert-authored ontology vocabulary</li>
        </ul>
      </div>
    </div>
    <div>
      <h3>Key Findings</h3>
      {callout('<strong>F1.</strong> Model size alone does not predict performance. Llama-3.1-8B (8B) outperforms Llama-2-70B (70B) on parse reliability.')}
      {callout('<strong>F2.</strong> Gemini 3.1 Pro is the only model with 100% parse success across all configurations — but also the highest hallucinator (0.623).','yellow')}
      {callout('<strong>F3.</strong> class_label_ratio = 1.0 for all models. property_label_ratio reveals hidden gaps — Llama-2-70B labels only 71.7% of properties.','yellow')}
      {callout('<strong>F4.</strong> scenario-cq-reasoning achieves the lowest hallucination (0.461 avg) and best scenario coverage — best faithfulness-quality trade-off.','green')}
      {callout('<strong>F5.</strong> BLOOMZ-7B1: 0% parse rate across all configs. Instruction-tuned BLOOM-family models still cannot produce valid OWL Turtle.','red')}
      {callout('<strong>F6.</strong> Combined F1 is universally low (0.06–0.13) — LLMs diverge in naming conventions from expert ontologists even when semantically correct.')}
    </div>
  </div>
  {callout('<strong>Recommendations:</strong> Proprietary — <strong>Gemini 3.1 Pro + scenario-only</strong> (structural winner) or <strong>GPT-5.4 + scenario-cq-constraints</strong> (functional runner-up). Open-source — <strong>Llama-3.1-8B + scenario-only</strong> (structural track) and <strong>Llama-2-70B + cq-only</strong> (functional track, score 0.534). The two open-source configurations are selected for expert human evaluation.','green')}
</div>

<div style="text-align:center;padding:40px;color:var(--muted);font-size:12px;border-top:1px solid var(--border);">
  ODPGen Research Presentation &nbsp;·&nbsp; {s['n_models_all']} Models &nbsp;·&nbsp;
  {s['n_configs']} Configs &nbsp;·&nbsp; {s['n_scenarios']} Scenarios &nbsp;·&nbsp;
  {s['n_runs']} Runs &nbsp;·&nbsp; {s['n_metrics']} Metrics
</div>
</body></html>"""

    HTML_OUT.write_text(html, encoding="utf-8")
    print(f"Saved: {HTML_OUT}")


# ═══════════════════════════════════════════════════════════════════════════════
# PPTX GENERATION
# ═══════════════════════════════════════════════════════════════════════════════

from pptx import Presentation as _Prs
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Palette
BG       = RGBColor(0x0f,0x11,0x17); SURFACE  = RGBColor(0x1a,0x1d,0x2e)
SURFACE2 = RGBColor(0x23,0x27,0x40); ACCENT   = RGBColor(0x6c,0x8e,0xbf)
ACCENT2  = RGBColor(0x8a,0xb4,0xf8); GREEN    = RGBColor(0x4c,0xaf,0x7d)
YELLOW   = RGBColor(0xf0,0xc0,0x60); RED      = RGBColor(0xe0,0x6c,0x6c)
PURPLE   = RGBColor(0xb3,0x9d,0xdb); WHITE    = RGBColor(0xe8,0xea,0xf0)
MUTED    = RGBColor(0x88,0x92,0xa4); DARK_BL  = RGBColor(0x12,0x18,0x30)
W = Inches(13.33); H = Inches(7.5)


def _rgb(c): return RGBColor(*c) if isinstance(c, tuple) else c
def _dim(c, d=8): return RGBColor(c[0]//d, c[1]//d, c[2]//d)
def _dim3(c, d=3): return RGBColor(c[0]//d, c[1]//d, c[2]//d)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color


def _box(slide, x, y, w, h, fill=None, border=None, bpt=1.5):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.line.width = Pt(bpt)
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = _rgb(fill)
    else: sh.fill.background()
    if border: sh.line.color.rgb = _rgb(border)
    else: sh.line.fill.background()
    return sh


def _tx(slide, x, y, w, h, text, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = str(text)
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb = _rgb(color)
    return tb


def _header(slide, num, title, y=Inches(0.55)):
    _box(slide, Inches(0.5), y, Inches(0.55), Inches(0.35), SURFACE2)
    _tx(slide, Inches(0.5), y+Pt(2), Inches(0.55), Inches(0.35), num, 11, color=MUTED, align=PP_ALIGN.CENTER)
    _tx(slide, Inches(1.2), y-Pt(2), Inches(11.5), Inches(0.45), title, 26, bold=True, color=ACCENT2)
    ln = _box(slide, Inches(0.5), y+Inches(0.43), Inches(12.3), Pt(1), SURFACE2)
    ln.line.fill.background()


def _stat(slide, x, y, w, h, num, label, nc=ACCENT2):
    _box(slide, x, y, w, h, SURFACE, SURFACE2)
    _tx(slide, x, y+Inches(0.18), w, Inches(0.7), num, 42, bold=True, color=nc, align=PP_ALIGN.CENTER)
    _tx(slide, x, y+Inches(0.85), w, Inches(0.35), label, 12, color=MUTED, align=PP_ALIGN.CENTER)


def _card(slide, x, y, w, h, title, bullets, tc=YELLOW, bc=SURFACE2):
    _box(slide, x, y, w, h, SURFACE, _rgb(bc))
    _tx(slide, x+Pt(6), y+Pt(7), w-Pt(12), Inches(0.3), title, 12, bold=True, color=_rgb(tc))
    body = "\n".join(f"• {b}" for b in bullets)
    _tx(slide, x+Pt(6), y+Inches(0.38), w-Pt(12), h-Inches(0.45), body, 10, color=WHITE)


def _callout(slide, x, y, w, h, text, ac=ACCENT):
    ac = _rgb(ac)
    bg_c = RGBColor(ac[0]//8, ac[1]//8, ac[2]//8)
    _box(slide, x, y, w, h, bg_c, ac, 2)
    _tx(slide, x+Inches(0.15), y+Pt(6), w-Inches(0.3), h-Pt(12), text, 11, color=WHITE)


def _hbar(slide, x, y, w, h, pct, label, vtext, col=GREEN):
    _tx(slide, x, y+Pt(3), Inches(2.3), h, label, 11, color=WHITE, align=PP_ALIGN.RIGHT)
    track = _box(slide, x+Inches(2.4), y, w, h, SURFACE2)
    track.line.fill.background()
    fw = max(Pt(2), int(w * max(pct, 0.01)))
    bar = _box(slide, x+Inches(2.4), y, fw, h, _rgb(col))
    bar.line.fill.background()
    _tx(slide, x+Inches(2.4)+fw+Pt(4), y+Pt(3), Inches(1.5), h, vtext, 10, bold=True, color=WHITE)


def _model_color(model_id):
    c = _m(model_id, "pptx_color", (0x8a, 0xb4, 0xf8))
    return _rgb(c) if isinstance(c, tuple) else c


def _parse_bar_color(rate):
    if rate >= 0.7: return GREEN
    if rate > 0:    return YELLOW
    return RED


def build_pptx(s: dict) -> None:
    prs = _Prs()
    prs.slide_width = W; prs.slide_height = H
    ps, qual, cfg_qual = s["ps"], s["qual"], s["cfg_qual"]
    models   = s["models_ordered"]
    configs  = s["configs_present"]
    cross    = s["cross"]
    avail_q  = s["avail_q"]

    # ─── SLIDE 1 — TITLE ──────────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl)
    lp = _box(sl, 0, 0, Inches(6.5), H, DARK_BL); lp.line.fill.background()
    ab = _box(sl, 0, 0, Inches(0.08), H, ACCENT);  ab.line.fill.background()
    _tx(sl, Inches(0.5), Inches(0.9), Inches(3.5), Inches(0.35), "RESEARCH PRESENTATION", 10, color=ACCENT)
    _tx(sl, Inches(0.5), Inches(1.35), Inches(5.8), Inches(1.3), "ODPGen", 70, bold=True, color=WHITE)
    _tx(sl, Inches(0.5), Inches(2.78), Inches(5.8), Inches(0.8),
        "Benchmarking Large Language Models on\nOntology Design Pattern Generation", 17, color=MUTED, italic=True)
    d = _box(sl, Inches(0.5), Inches(3.68), Inches(5.5), Pt(1.5), ACCENT); d.line.fill.background()
    meta = [("Models", str(s["n_models_all"])), ("Prompt Configs", str(s["n_configs"])),
            ("ODP Scenarios", str(s["n_scenarios"])), ("Total Runs", str(s["n_runs"]))]
    for i,(lbl,val) in enumerate(meta):
        col=i%2; row=i//2
        _tx(sl, Inches(0.5+col*2.8), Inches(3.85+row*0.85), Inches(2.6), Inches(0.28), lbl, 10, color=MUTED)
        _tx(sl, Inches(0.5+col*2.8), Inches(4.13+row*0.85), Inches(2.6), Inches(0.4),  val, 18, bold=True, color=ACCENT2)
    for i,(num,lbl) in enumerate([(str(s["n_runs"]),"runs"),(str(s["n_scenarios"]),"scenarios"),
                                   (str(s["n_configs"]),"configs"),(str(s["n_metrics"]),"metrics")]):
        rx=Inches(7.2+(i%2)*2.85); ry=Inches(1.6+(i//2)*2.4)
        _stat(sl, rx, ry, Inches(2.5), Inches(1.9), num, lbl)

    # ─── SLIDE 2 — OVERVIEW ───────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"01","Overview & Research Questions")
    _card(sl, Inches(0.5), Inches(1.22), Inches(5.5), Inches(3.0), "Research Questions",
          ["RQ1  How well can LLMs generate valid OWL/Turtle from scenarios?",
           "RQ2  How logically consistent and error-free are outputs?",
           "RQ3  What is the impact of prompting strategy?",
           "RQ4  How faithful are outputs to the ground-truth scenario?",
           "RQ5  Can models generalize to recently-published ODPs?"], ACCENT2, ACCENT)
    stats_data = [(str(s["n_models_all"]),"LLMs"),(str(s["n_configs"]),"Configs"),
                  (str(s["n_scenarios"]),"Scenarios"),(str(s["n_runs"]),"Runs")]
    for i,(num,lbl) in enumerate(stats_data):
        _stat(sl, Inches(6.3+i*1.75), Inches(1.22), Inches(1.6), Inches(1.4), num, lbl)
    _callout(sl, Inches(0.5), Inches(4.35), Inches(12.3), Inches(0.65),
             "Core hypothesis: structured prompting (CQs + reasoning or self-validation) should improve "
             "faithfulness — but may trade off syntactic correctness for weaker models.")
    # mini stat row
    for i,(num,lbl,nc) in enumerate([(str(s["n_runs"]),"evaluated outputs",ACCENT2),
                                      (str(s["n_scenarios"]),"ODP scenarios",GREEN),
                                      (str(s["n_configs"]),"prompt configs",YELLOW),
                                      (str(s["n_metrics"]),"metrics",PURPLE)]):
        _stat(sl, Inches(0.5+i*3.1), Inches(5.2), Inches(2.8), Inches(1.9), num, lbl, _rgb(nc))

    # ─── SLIDE 2b — DATASET ───────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"02","Dataset — 14 ODP Scenarios")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),
        "Scenarios sourced from published ODP papers (2023–2026), each with scenario text, Competency Questions, and a ground-truth OWL ontology.",
        12,color=MUTED)
    ds_xs=[Inches(0.5),Inches(2.55),Inches(7.7),Inches(8.7),Inches(10.1)]
    ds_ws=[Inches(2.0),Inches(5.1),Inches(0.95),Inches(1.35),Inches(2.5)]
    ds_hdrs=["Scenario ID","Domain","# CQs","Year","Unseen Status"]
    for hx,hw,hl in zip(ds_xs,ds_ws,ds_hdrs):
        _box(sl,hx,Inches(1.55),hw,Inches(0.3),SURFACE2)
        _tx(sl,hx+Pt(3),Inches(1.58),hw,Inches(0.26),hl,8,bold=True,color=ACCENT2)
    row_h=Inches(0.37)
    for i,sc in enumerate(s["scenarios"]):
        ry=Inches(1.88)+i*row_h
        rbg=SURFACE if i%2==0 else BG
        _box(sl,Inches(0.5),ry,Inches(12.1),row_h,rbg)
        yr=sc["year"]
        if yr<=2023: utag="Likely Seen"; uc=GREEN
        elif yr==2024: utag="Borderline"; uc=YELLOW
        else: utag="Plausible Unseen"; uc=RED
        vals=[sc["id"],sc.get("domain","—"),str(sc["cqs"]),str(yr),utag]
        vcols=[ACCENT2,WHITE,WHITE,MUTED,uc]
        bolds=[True,False,False,False,False]
        for hx,hw,val,vc,bold in zip(ds_xs,ds_ws,vals,vcols,bolds):
            _tx(sl,hx+Pt(3),ry+Pt(5),hw-Pt(4),row_h-Pt(4),val,8,bold=bold,color=_rgb(vc))
    _callout(sl,Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.38),
             "CQ range: 3 (Role-Dependent Names) to 47 (AISHIP Maritime). 9 of 14 scenarios plausible-unseen for all evaluated models.",YELLOW)

    # ─── SLIDE 2c — WORKED EXAMPLE ────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"03","Worked Example — Scenario 2023-133-01: Causal Bayesian Network")
    # Left: scenario + CQs + ground truth
    _box(sl,Inches(0.5),Inches(1.22),Inches(6.0),Inches(2.05),RGBColor(0x08,0x0e,0x1a),ACCENT,1.5)
    _tx(sl,Inches(0.65),Inches(1.32),Inches(5.7),Inches(0.28),"Scenario Text",11,bold=True,color=ACCENT2)
    _tx(sl,Inches(0.65),Inches(1.62),Inches(5.7),Inches(1.55),
        "Sprinkler use case (CBN). Five event nodes: rainy season, sprinkler on, raining, pavement wet, pavement slippery. "
        "Directed edges carry numeric effect weights (e.g., rainy season->rain = 0.6). "
        "Rain->slippery is mediated by wet pavement.",
        9,color=WHITE)
    _box(sl,Inches(0.5),Inches(3.32),Inches(6.0),Inches(1.65),RGBColor(0x08,0x0e,0x1a),SURFACE2,1)
    _tx(sl,Inches(0.65),Inches(3.42),Inches(5.7),Inches(0.28),"4 Competency Questions",11,bold=True,color=ACCENT2)
    cq_texts=["1. What events caused the pavement to be slippery?",
              "2. Outcome of turning the sprinkler on/off?",
              "3. Strength of causal link: sprinkler -> wet pavement?",
              "4. If pavement not wet, effect on slipperiness? (counterfactual)"]
    for k,cq in enumerate(cq_texts):
        _tx(sl,Inches(0.65),Inches(3.72)+k*Inches(0.35),Inches(5.7),Inches(0.32),cq,9,color=WHITE)
    _box(sl,Inches(0.5),Inches(5.02),Inches(6.0),Inches(1.65),RGBColor(0x04,0x12,0x08),GREEN,1.5)
    _tx(sl,Inches(0.65),Inches(5.12),Inches(5.7),Inches(0.28),"Ground-Truth ODP (Expert-Authored)",11,bold=True,color=GREEN)
    gt_pts=["5 classes: Event, AbstractEvent, ConcreteEvent, Causes, EffectWeight",
            "Causes is a CLASS (reification) -- not a property",
            "hasTreatment, hasOutcome, hasMediator (all on Causes->Event)",
            "hasEffectWeight (Causes->EffectWeight)",
            "Abstract, reusable -- no domain-specific sprinkler/rain classes"]
    for k,pt in enumerate(gt_pts):
        _tx(sl,Inches(0.65),Inches(5.42)+k*Inches(0.24),Inches(5.7),Inches(0.22),f"• {pt}",8,color=WHITE)
    # Right: generated outputs comparison
    _tx(sl,Inches(6.7),Inches(1.22),Inches(6.1),Inches(0.28),"Generated Outputs — scenario-only",13,bold=True,color=ACCENT2)
    _box(sl,Inches(6.7),Inches(1.55),Inches(6.1),Inches(2.0),RGBColor(0x08,0x0e,0x1a),RGBColor(0x2a,0x3e,0x6e),1)
    _tx(sl,Inches(6.85),Inches(1.62),Inches(3.0),Inches(0.28),"Gemini 3.1 Pro",11,bold=True,color=ACCENT2)
    _tx(sl,Inches(9.9),Inches(1.62),Inches(2.0),Inches(0.28),"Parse: SUCCESS",10,bold=True,color=GREEN,align=PP_ALIGN.RIGHT)
    gemini_code=(":Event a owl:Class .\n"
                 ":CausalRelation a owl:Class ; owl:disjointWith :Event .\n"
                 ":RainySeasonEvent rdfs:subClassOf :Event .\n"
                 ":SprinklerOnEvent rdfs:subClassOf :Event .\n"
                 ":PavementWetEvent rdfs:subClassOf :Event .\n"
                 ":PavementSlipperyEvent rdfs:subClassOf :Event .\n"
                 ":hasSourceEvent, :hasTargetEvent, :isMediatedBy\n"
                 "  rdfs:domain :CausalRelation; rdfs:range :Event .\n"
                 ":hasEffectWeight a owl:DatatypeProperty; rdfs:range xsd:decimal .")
    _tx(sl,Inches(6.85),Inches(1.95),Inches(5.85),Inches(1.25),gemini_code,7,color=RGBColor(0xc3,0xc8,0xd8))
    _tx(sl,Inches(6.85),Inches(3.18),Inches(5.85),Inches(0.3),"+ 5 domain-specific subclasses (scenario-grounded, not in GT)",8,color=YELLOW,italic=True)
    _box(sl,Inches(6.7),Inches(3.55),Inches(6.1),Inches(1.9),RGBColor(0x04,0x12,0x08),RGBColor(0x1e,0x5e,0x30),1)
    _tx(sl,Inches(6.85),Inches(3.62),Inches(3.0),Inches(0.28),"Llama-3.1-8B",11,bold=True,color=GREEN)
    _tx(sl,Inches(9.9),Inches(3.62),Inches(2.0),Inches(0.28),"Parse: SUCCESS",10,bold=True,color=GREEN,align=PP_ALIGN.RIGHT)
    llama_code=(":WeatherEvent, :RainySeason, :SprinklerStatus\n"
                ":SprinklerOn rdfs:subClassOf :SprinklerStatus\n"
                ":Weather, :Rain rdfs:subClassOf :Weather\n"
                ":PavementCondition\n"
                ":PavementWet, :PavementSlippery\n"
                "  rdfs:subClassOf :PavementCondition\n"
                ":CausalRelation, :EffectWeight (DatatypeProperty)\n"
                "-> more domain-specific hierarchy; more superfluous")
    _tx(sl,Inches(6.85),Inches(3.95),Inches(5.85),Inches(1.15),llama_code,7,color=RGBColor(0xc3,0xc8,0xd8))
    _callout(sl,Inches(6.7),Inches(5.52),Inches(6.1),Inches(1.1),
             "Systematic finding: both models literalize the scenario into domain-specific class hierarchies. "
             "The expert ODP uses reification (Causes as class) for a domain-agnostic, reusable pattern. "
             "LLMs struggle to generate abstract patterns from narrative descriptions.",YELLOW)

    # ─── SLIDE 2d — PROMPT TEMPLATES FULL TEXT ────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"04","Prompt Templates — Actual Text Sent to Models")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),
        "All 5 configurations share the same OWL/Turtle output contract. They differ in what input they receive and what output structure they require.",
        12,color=MUTED)
    # Prompt 1: scenario-only (left)
    _box(sl,Inches(0.5),Inches(1.55),Inches(6.0),Inches(5.7),SURFACE,ACCENT,1.5)
    _tx(sl,Inches(0.65),Inches(1.65),Inches(5.7),Inches(0.28),"Config 1 — scenario-only",12,bold=True,color=ACCENT2)
    p1_text=("You are an ontology engineer. Generate\n"
             "an OWL ODP in Turtle syntax from the\n"
             "scenario below.\n\n"
             "Requirements:\n"
             "- Stay strictly within scenario scope.\n"
             "- Use clear, self-explanatory names.\n"
             "- Include minimal but sufficient axioms.\n"
             "- Declare every class with `a owl:Class`,\n"
             "  every property with its OWL type.\n"
             "- Do NOT create individual instances.\n"
             "- Add rdfs:label + rdfs:comment to all.\n"
             "- Declare owl:Ontology header.\n\n"
             "@prefix owl: <...owl#> .\n"
             "@prefix rdfs: <...rdf-schema#> .\n"
             "@prefix : <http://example.org/odp#> .\n\n"
             "Output -- in this order:\n"
             "1. Turtle in ```turtle...``` block FIRST.\n"
             "2. Intent . Requirements . Assumptions.\n\n"
             "Scenario:\n{{SCENARIO_TEXT}}")
    _tx(sl,Inches(0.65),Inches(2.0),Inches(5.7),Inches(5.0),p1_text,8,color=RGBColor(0xc3,0xc8,0xd8))
    # Prompt 5: scenario-cq-reasoning (right)
    _box(sl,Inches(6.7),Inches(1.55),Inches(6.1),Inches(5.7),SURFACE,GREEN,1.5)
    _tx(sl,Inches(6.85),Inches(1.65),Inches(5.85),Inches(0.28),"Config 5 — scenario-cq-reasoning",12,bold=True,color=GREEN)
    p5_text=("You are an ontology engineer. Create\n"
             "an OWL ODP in Turtle syntax.\n\n"
             "Before writing, reason through these\n"
             "steps internally (do NOT write yet):\n"
             "1. Extract atomic requirements.\n"
             "2. Propose candidate classes/properties.\n"
             "3. Draft axioms (subclass, domain, range).\n"
             "4. Verify each CQ is answered by axioms.\n"
             "5. Refine axioms to cover uncovered CQs.\n"
             "6. Remove entities not traceable to reqs/CQs.\n\n"
             "Constraints:\n"
             "- schema axioms only; no instances\n"
             "- rdfs:label + rdfs:comment mandatory\n"
             "- owl:Ontology header mandatory\n"
             "- minimal, reusable; no extra facts\n\n"
             "Output -- in this order:\n"
             "1. Turtle in ```turtle...``` block FIRST.\n"
             "2. Reasoning summary: 1 bullet/step.\n"
             "3. Traceability: req/CQ -> axiom(s).\n\n"
             "Scenario: {{SCENARIO_TEXT}}\n"
             "CQs: {{CQ_LIST}}")
    _tx(sl,Inches(6.85),Inches(2.0),Inches(5.85),Inches(5.0),p5_text,8,color=RGBColor(0xc3,0xc8,0xd8))

    # ─── SLIDE 3 — MODELS ─────────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"02","Models Evaluated")
    col_hdrs = ["Model","Provider","Parameters","Training Cutoff","Parse Rate"]
    cxs=[Inches(0.5),Inches(2.6),Inches(5.3),Inches(7.0),Inches(9.2)]
    cws=[Inches(2.05),Inches(2.65),Inches(1.65),Inches(2.15),Inches(1.5)]
    for hx,hw,hl in zip(cxs,cws,col_hdrs):
        _box(sl,hx,Inches(1.22),hw,Inches(0.33),SURFACE2); _tx(sl,hx+Pt(4),Inches(1.25),hw,Inches(0.28),hl,9,bold=True,color=ACCENT2)
    for i,m in enumerate(models):
        ry=Inches(1.58)+i*Inches(0.77)
        mc=_model_color(m); bg_c=RGBColor(mc[0]//14,mc[1]//14,mc[2]//14)
        _box(sl,Inches(0.5),ry,Inches(10.75),Inches(0.72),bg_c,RGBColor(mc[0]//4,mc[1]//4,mc[2]//4),1)
        rate=ps.loc[m,"rate"] if m in ps.index else 0
        rate_str=_pct(rate); rate_col=_parse_bar_color(rate)
        vals=[_m(m,"name",m),_m(m,"provider","—"),_m(m,"params","—"),_m(m,"cutoff","—"),rate_str]
        cols2=[mc,MUTED,WHITE,MUTED,rate_col]; bolds=[True,False,False,False,True]
        for hx,hw,val,vc,bold in zip(cxs,cws,vals,cols2,bolds):
            _tx(sl,hx+Pt(6),ry+Pt(10),hw-Pt(6),Inches(0.55),val,11,bold=bold,color=_rgb(vc))
    _callout(sl,Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.38),
             "Temperature=0.0 for all runs. Open-source: HuggingFace Transformers. GPT-5.4: OpenAI API. Gemini: Google GenAI SDK.")

    # ─── SLIDE 4 — PROMPTS ────────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"03","Prompting Strategies")
    cfg_info = [
        ("1. scenario-only",["Only scenario text — no CQs","Model infers scope from narrative","Highest scenario vocabulary coverage"],ACCENT2),
        ("2. cq-only",["Only Competency Questions provided","Highest CQ coverage (0.858 avg)","Also highest hallucination (0.586)"],ACCENT2),
        ("3. scenario-cq  [Balanced]",["Both scenario + CQs","Requires scenario→axiom & CQ→axiom mapping","Good balance of CQ + scenario coverage"],GREEN),
        ("4. scenario-cq-constraints  [Validation]",["Adds 6-item self-validation checklist","Syntax, domain/range, coherence, scope checks","Lower parse rate — extra text confuses parsers"],YELLOW),
        ("5. scenario-cq-reasoning  [Faithfulness]",["6-step reasoning: extract → propose → draft → verify → remove → output","Lowest hallucination overall","Best scenario coverage among CQ-aware configs"],GREEN),
    ]
    positions = [(Inches(0.5),Inches(1.22),Inches(3.9),Inches(2.1)),
                 (Inches(4.55),Inches(1.22),Inches(3.9),Inches(2.1)),
                 (Inches(8.6),Inches(1.22),Inches(4.2),Inches(2.1)),
                 (Inches(0.5),Inches(3.48),Inches(5.9),Inches(2.1)),
                 (Inches(6.55),Inches(3.48),Inches(6.25),Inches(2.1))]
    for (title,bullets,tc),(cx,cy,cw,ch) in zip(cfg_info,positions):
        _card(sl,cx,cy,cw,ch,title,bullets,tc)
    _callout(sl,Inches(0.5),Inches(5.75),Inches(12.3),Inches(0.52),
             "All prompts use @prefix owl/rdf/rdfs/xsd and require owl:Ontology with rdfs:label and rdfs:comment.")

    # ─── SLIDE 5 — EVALUATION ─────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"04","Evaluation Methodology — 4 Layers")
    layers = [
        ("Layer 1 — Formal / Syntactic Quality",ACCENT,
         ["parse_success — Valid Turtle syntax?","triple_count — Total RDF triples",
          "ontology_declared — owl:Ontology header?","label_coverage_ratio — all entities with rdfs:label",
          "class_label_ratio [NEW] — owl:Class entities with rdfs:label",
          "property_label_ratio [NEW] — Properties with rdfs:label",
          "comment_coverage_ratio, domain_range_ratio"]),
        ("Layer 2 — Scenario Faithfulness",GREEN,
         ["cq_coverage_ratio — CQs answerable (token overlap)","cq_mean_token_recall — avg recall per CQ",
          "scenario_coverage_ratio — scenario vocab in ontology","hallucination_ratio — unsupported vocab fraction"]),
        ("Layer 3 — Semantic Similarity vs Ground Truth",YELLOW,
         ["class_name_F1 — name match vs ground-truth classes","property_name_F1 — same for properties",
          "combined_f1 — macro avg F1","matched_class_names — classes in both generated & GT",
          "matched_property_names — properties in both; e.g. 'event', 'has effect'"]),
        ("Layer 4 — Feasibility & Unseen Analysis",PURPLE,
         ["Each (model, ODP) labelled plausible_unseen or likely_seen",
          "Compares ODP pub date vs model training cutoff","Confidence: high / medium / low",
          "Future: human eval — Correctness, Completeness, Clarity (1–5)"]),
    ]
    for i,(title,color,bullets) in enumerate(layers):
        col=i%2; row=i//2
        cx=Inches(0.5+col*6.45); cy=Inches(1.22+row*2.85)
        _card(sl,cx,cy,Inches(6.1),Inches(2.7),title,bullets,_rgb(color),_dim3(color))
    _callout(sl,Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.38),
             f"Pipeline → results/summary.csv  ({s['n_runs']} rows × {s['n_metrics']} columns, incl. class_label_ratio & property_label_ratio)")

    # ─── SLIDE 6 — RAW OUTPUTS ────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"05","Raw Output Examples — Scenario 2023-133-01")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.3),
        'Scenario: "Model a causal Bayesian network over events with weighted causal edges as an ODP."',
        11,color=MUTED,italic=True)
    cols3 = 3; col_w = Inches(4.2); col_gap = Inches(0.17)
    for i,m in enumerate(models[:6]):
        col=i%cols3; row=i//cols3
        cx=Inches(0.5)+col*(col_w+col_gap); cy=Inches(1.62)+row*Inches(2.6)
        mc=_model_color(m); bg_c=RGBColor(mc[0]//12,mc[1]//12,mc[2]//12)
        rate=ps.loc[m,"rate"] if m in ps.index else 0
        _box(sl,cx,cy,col_w,Inches(2.45),bg_c,RGBColor(mc[0]//4,mc[1]//4,mc[2]//4),1)
        status="SUCCESS" if rate>0 else "FAIL"; sc=GREEN if rate>0 else RED
        _tx(sl,cx+Pt(6),cy+Pt(6),col_w-Pt(60),Inches(0.28),_m(m,"name",m),12,bold=True,color=_rgb(mc))
        _tx(sl,cx+col_w-Inches(0.95),cy+Pt(8),Inches(0.88),Inches(0.24),status,9,bold=True,color=_rgb(sc),align=PP_ALIGN.CENTER)
        code_bg=_box(sl,cx+Pt(5),cy+Inches(0.38),col_w-Pt(10),Inches(1.65),SURFACE2); code_bg.line.fill.background()
        snippet=(s["sample_outputs"].get(m,"(no output)")[:350]).replace("\n"," ↩ ")
        _tx(sl,cx+Pt(8),cy+Inches(0.41),col_w-Pt(16),Inches(1.6),snippet,7,color=RGBColor(0xc3,0xc8,0xd8))

    # ─── SLIDE 7 — PARSE SUCCESS ──────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"06","Results — Parse Success Rate")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),
        f"Parse success is the primary gate metric. Total: {s['n_runs']} runs.",12,color=MUTED)
    _tx(sl,Inches(0.5),Inches(1.55),Inches(6.0),Inches(0.3),"By Model",14,bold=True,color=ACCENT2)
    for i,m in enumerate(models):
        rate=ps.loc[m,"rate"] if m in ps.index else 0
        parsed=int(ps.loc[m,"parsed"]) if m in ps.index else 0
        total=int(ps.loc[m,"total"])  if m in ps.index else 0
        _hbar(sl,Inches(0.5),Inches(1.9)+i*Inches(0.62),Inches(5.8),Inches(0.45),
              rate,_m(m,"name",m),f"{_pct(rate)}  ({parsed}/{total})",_parse_bar_color(rate))
    _tx(sl,Inches(7.0),Inches(1.55),Inches(5.8),Inches(0.3),"By Configuration",14,bold=True,color=ACCENT2)
    cfg_ps=(s["df"].groupby("config")["parse_success"].agg(["sum","count","mean"])
              .rename(columns={"mean":"rate"}))
    for i,c in enumerate(configs):
        if c not in cfg_ps.index: continue
        r=cfg_ps.loc[c,"rate"]
        _hbar(sl,Inches(7.0),Inches(1.9)+i*Inches(0.62),Inches(5.8),Inches(0.45),
              r,c,f"{_pct(r)}  ({int(cfg_ps.loc[c,'sum'])}/{int(cfg_ps.loc[c,'count'])})",ACCENT2)
    _callout(sl,Inches(0.5),Inches(6.45),Inches(12.3),Inches(0.52),
             "Gemini 3.1 Pro: 100% parse success across ALL configs — the only model to achieve this. "
             "GPT-5.4: 86% on cq-only but 0% on reasoning/constraints.",GREEN)

    # ─── SLIDE 8 — QUALITY METRICS ────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"07","Results — Quality Metrics (Successful Parses)")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),
        f"Metrics computed only when parse_success=True.  {s['n_runs']} runs × {s['n_metrics']} metrics.",11,color=MUTED)
    show_q=["cq_coverage_ratio","scenario_coverage_ratio","hallucination_ratio",
            "class_label_ratio","property_label_ratio","combined_f1"]
    show_q=[c for c in show_q if c in avail_q]
    q_short={"cq_coverage_ratio":"CQ Cov.","scenario_coverage_ratio":"Scen. Cov.",
              "hallucination_ratio":"Halluc.","class_label_ratio":"Cls Lbl★",
              "property_label_ratio":"Prop Lbl★","combined_f1":"F1"}
    hdr_lbls=["Model","Parse"]+[q_short.get(c,c) for c in show_q]
    n_cols=len(hdr_lbls)
    col_w_each=Inches(12.3)/n_cols
    hdr_xs=[Inches(0.5)+i*col_w_each for i in range(n_cols)]
    for hx,hl in zip(hdr_xs,hdr_lbls):
        _box(sl,hx,Inches(1.55),col_w_each,Inches(0.3),SURFACE2); _tx(sl,hx+Pt(3),Inches(1.58),col_w_each,Inches(0.26),hl,8,bold=True,color=ACCENT2)
    for i,m in enumerate(models):
        ry=Inches(1.88)+i*Inches(0.52)
        mc=_model_color(m); bg_c=RGBColor(mc[0]//14,mc[1]//14,mc[2]//14)
        _box(sl,Inches(0.5),ry,Inches(12.3),Inches(0.49),bg_c);
        rate=ps.loc[m,"rate"] if m in ps.index else 0
        rc=_parse_bar_color(rate)
        vals=[_m(m,"name",m),_pct(rate)]
        vcols=[_rgb(mc),_rgb(rc)]
        bolds=[True,True]
        if not qual.empty and m in qual.index:
            for c in show_q:
                v=_fmt(qual.loc[m,c]) if c in qual.columns else "—"
                # colour by rank
                if c in qual.columns and not qual[c].isna().all():
                    rk=qual[c].rank(ascending=(c!="hallucination_ratio"),pct=True).get(m,0.5)
                    vc=GREEN if rk>=0.75 else(RED if rk<=0.25 else WHITE)
                else: vc=WHITE
                vals.append(v); vcols.append(_rgb(vc)); bolds.append(False)
        else:
            for _ in show_q: vals.append("—"); vcols.append(MUTED); bolds.append(False)
        for hx,val,vc,bold in zip(hdr_xs,vals,vcols,bolds):
            _tx(sl,hx+Pt(3),ry+Pt(6),col_w_each-Pt(3),Inches(0.4),val,10,bold=bold,color=vc)
    _callout(sl,Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.38),
             "class_label_ratio=1.000 for ALL models. property_label_ratio reveals gap: Llama-2-70B=0.717(worst), Mistral=0.966(best).",YELLOW)

    # ─── SLIDE 9 — HEATMAP ────────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"08","Model × Config Parse Success Heatmap")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),"Parse success rate per (model, config) pair. 14 scenarios each.",12,color=MUTED)
    n_cfg=len(configs); cell_w=Inches(10.0)/n_cfg; cell_h=Inches(0.65); start_x=Inches(2.65); start_y=Inches(1.55)
    for j,c in enumerate(configs+["Average"]):
        cx=start_x+j*cell_w
        _box(sl,cx,start_y,cell_w,Inches(0.48),SURFACE2); _tx(sl,cx+Pt(2),start_y+Pt(4),cell_w-Pt(4),Inches(0.4),c if c!="Average" else "Avg",8,bold=True,color=ACCENT2,align=PP_ALIGN.CENTER)
    for i,m in enumerate(models):
        ry=start_y+Inches(0.5)+i*cell_h
        mc=_model_color(m); bg_c=RGBColor(mc[0]//14,mc[1]//14,mc[2]//14)
        _box(sl,Inches(0.5),ry,Inches(2.1),cell_h,bg_c); _tx(sl,Inches(0.52),ry+Pt(10),Inches(2.05),cell_h-Pt(10),_m(m,"name",m),10,bold=True,color=_rgb(mc))
        for j,c in enumerate(configs):
            cx=start_x+j*cell_w
            if (m,c) in cross.index:
                r=cross.loc[(m,c),"rate"]; vc=_parse_bar_color(r)
                cbg=RGBColor(vc[0]//10,vc[1]//10,vc[2]//10)
                _box(sl,cx,ry,cell_w,cell_h,cbg,RGBColor(vc[0]//4,vc[1]//4,vc[2]//4),1)
                _tx(sl,cx,ry+Pt(12),cell_w,cell_h-Pt(12),_pct(r),14,bold=True,color=_rgb(vc),align=PP_ALIGN.CENTER)
            else:
                _box(sl,cx,ry,cell_w,cell_h,SURFACE,SURFACE2,1)
                _tx(sl,cx,ry+Pt(14),cell_w,Inches(0.3),"—",12,color=MUTED,align=PP_ALIGN.CENTER)
        # avg col
        avg_cx=start_x+n_cfg*cell_w; avg_r=ps.loc[m,"rate"] if m in ps.index else 0
        avc=_parse_bar_color(avg_r); abg=RGBColor(avc[0]//10,avc[1]//10,avc[2]//10)
        _box(sl,avg_cx,ry,cell_w,cell_h,abg,RGBColor(avc[0]//4,avc[1]//4,avc[2]//4),2)
        _tx(sl,avg_cx,ry+Pt(12),cell_w,cell_h-Pt(12),_pct(avg_r),14,bold=True,color=_rgb(avc),align=PP_ALIGN.CENTER)
    _callout(sl,Inches(0.5),Inches(6.88),Inches(7.5),Inches(0.52),
             "Gemini 3.1 Pro: 100% across every config — unique. BLOOMZ-7B1: 0% across every config.",GREEN)
    _callout(sl,Inches(8.2),Inches(6.88),Inches(4.6),Inches(0.52),
             "GPT-5.4: 86% on cq-only, 0% on constraints/reasoning. Prompts break Turtle extraction.",YELLOW)

    # ─── SLIDE 9b — STRUCTURAL RANKINGS ──────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"09b","Pilot Results — Structural Evaluation Rankings")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),
        "Structural score = 0.5 × parse_rate + 0.5 × avg_OOPS_score. "
        "OOPS score: freedom from common OWL pitfalls (0=worst, 1=best). Computed per (model, config) pair, 14 scenarios.",
        11,color=MUTED)
    # ── Proprietary table (left) ───────────────────────────────────────────────
    _tx(sl,Inches(0.5),Inches(1.55),Inches(6.1),Inches(0.28),"Proprietary Models — Top 4 Configurations",13,bold=True,color=ACCENT2)
    prop_hdrs=["Rank","Model","Config","Struct. Score","OOPS","Parse Rate"]
    prop_xs=[Inches(0.5),Inches(0.98),Inches(2.55),Inches(4.78),Inches(5.65),Inches(6.3)]
    prop_ws=[Inches(0.45),Inches(1.52),Inches(2.18),Inches(0.82),Inches(0.62),Inches(0.6)]
    for hx,hw,hl in zip(prop_xs,prop_ws,prop_hdrs):
        _box(sl,hx,Inches(1.88),hw,Inches(0.28),SURFACE2)
        _tx(sl,hx+Pt(2),Inches(1.91),hw,Inches(0.24),hl,7,bold=True,color=ACCENT2)
    prop_rows=[
        ("1","Gemini 3.1 Pro","scenario-only","0.982","0.964","100%",True),
        ("2","Gemini 3.1 Pro","scenario-cq-reasoning","0.857","0.714","100%",False),
        ("3","GPT-5.4","scenario-only","0.464","0.429","50%",False),
        ("4","GPT-5.4","scenario-cq-constraints","0.293","0.214","37%",False),
    ]
    for i,(rank,model,cfg,ss,oops,pr,best) in enumerate(prop_rows):
        ry=Inches(2.19)+i*Inches(0.52)
        rbg=RGBColor(0x07,0x16,0x0e) if best else (SURFACE if i%2==0 else BG)
        rborder=RGBColor(0x2a,0x6e,0x4a) if best else SURFACE2
        _box(sl,Inches(0.5),ry,Inches(6.6),Inches(0.49),rbg,rborder,1.5 if best else 0.5)
        ss_f=float(ss); oops_f=float(oops)
        ss_c=GREEN if ss_f>=0.7 else (YELLOW if ss_f>=0.4 else RED)
        mc=ACCENT2 if "Gemini" in model else PURPLE
        for hx,hw,val,vc,bold in zip(prop_xs,prop_ws,
            [rank,model,cfg,ss,oops,pr],
            [WHITE,_rgb(mc),MUTED,_rgb(ss_c),WHITE,WHITE],
            [True,True,False,True,False,False]):
            _tx(sl,hx+Pt(2),ry+Pt(6),hw-Pt(4),Inches(0.4),val,9,bold=bold,color=vc)
    _callout(sl,Inches(0.5),Inches(4.42),Inches(6.6),Inches(0.52),
             "Gemini scenario-only (0.982) = top structural result across ALL models and configs. "
             "Adding CQs does not always improve structural quality.",GREEN)
    # ── Open-source table (right) ─────────────────────────────────────────────
    _tx(sl,Inches(7.0),Inches(1.55),Inches(5.8),Inches(0.28),"Open-Source Models — Top 3 Configurations",13,bold=True,color=GREEN)
    os_hdrs=["Rank","Model","Config","Struct.","OOPS","Parse"]
    os_xs=[Inches(7.0),Inches(7.45),Inches(8.85),Inches(10.7),Inches(11.42),Inches(12.08)]
    os_ws=[Inches(0.42),Inches(1.35),Inches(1.82),Inches(0.68),Inches(0.62),Inches(0.58)]
    for hx,hw,hl in zip(os_xs,os_ws,os_hdrs):
        _box(sl,hx,Inches(1.88),hw,Inches(0.28),SURFACE2)
        _tx(sl,hx+Pt(2),Inches(1.91),hw,Inches(0.24),hl,7,bold=True,color=ACCENT2)
    os_rows=[
        ("1","Llama-3.1-8B","scenario-cq","0.616","0.375","85.7%",False),
        ("2★","Llama-3.1-8B","scenario-only","0.595","0.405","78.6%",True),
        ("3★","Llama-2-70B","cq-only","0.589","0.321","85.7%",True),
    ]
    for i,(rank,model,cfg,ss,oops,pr,selected) in enumerate(os_rows):
        ry=Inches(2.19)+i*Inches(0.52)
        rbg=RGBColor(0x04,0x12,0x08) if selected else (SURFACE if i%2==0 else BG)
        rborder=RGBColor(0x2a,0x6e,0x4a) if selected else SURFACE2
        _box(sl,Inches(7.0),ry,Inches(5.75),Inches(0.49),rbg,rborder,1.5 if selected else 0.5)
        ss_f=float(ss)
        ss_c=GREEN if ss_f>=0.7 else (YELLOW if ss_f>=0.4 else RED)
        mc=GREEN if "8B" in model else YELLOW
        for hx,hw,val,vc,bold in zip(os_xs,os_ws,
            [rank,model,cfg,ss,oops,pr],
            [WHITE,_rgb(mc),MUTED,_rgb(ss_c),WHITE,WHITE],
            [True,True,False,True,False,False]):
            _tx(sl,hx+Pt(2),ry+Pt(6),hw-Pt(4),Inches(0.4),val,9,bold=bold,color=vc)
    _callout(sl,Inches(7.0),Inches(3.42),Inches(5.75),Inches(0.52),
             "★ Selected for human evaluation. Scenario-only chosen over scenario-cq to test prompt minimalism.",YELLOW)
    _callout(sl,Inches(7.0),Inches(4.02),Inches(5.75),Inches(0.52),
             "Anomaly: Llama-3.1-8B produced 22 repeated identical triple blocks on scenario 2025-151-01. "
             "Parseable — retained as edge case.",RED)
    _callout(sl,Inches(0.5),Inches(5.05),Inches(12.3),Inches(0.45),
             "Best open-source structural score (0.616) is 37% below best proprietary (0.982) — "
             "a significant gap that human evaluation will further contextualize.",ACCENT)

    # ─── SLIDE 10 — FUNCTIONAL EVALUATION ────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"10","Functional Evaluation — Utility-Oriented Pilot Assessment")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.3),
        "Functional evaluation measures whether generated ODPs answer CQs without adding unsupported ontological commitments.",
        12,color=MUTED)
    # formula card (left)
    _card(sl,Inches(0.5),Inches(1.58),Inches(5.5),Inches(3.6),"Scoring Methodology",
          ["functional_score = 0.5 × manual_score + 0.5 × cq_pass_rate",
           "numb: total elements beyond ground truth",
           "numbnoscen: extra with no scenario grounding",
           "superfluous_penalty = numbnoscen×1.0 + (numb−numbnoscen)×0.3",
           "manual_score = 1 / (1 + 0.5 × superfluous_penalty)",
           "cq_pass_rate: fraction of CQs answered by the generated ODP",
           "Unsupported additions penalized ×1.0; scenario-justified ×0.3"],
          ACCENT2, ACCENT)
    # results table (right)
    _tx(sl,Inches(6.2),Inches(1.55),Inches(6.6),Inches(0.28),"Pilot Results — Best Config per Model Family",13,bold=True,color=ACCENT2)
    func_hdrs=["Model","Config","Avg Extra","Manual Score","CQ Pass","Func. Score"]
    func_xs=[Inches(6.2),Inches(7.5),Inches(9.0),Inches(10.0),Inches(11.1),Inches(12.0)]
    func_ws=[Inches(1.25),Inches(1.45),Inches(0.95),Inches(1.05),Inches(0.85),Inches(1.0)]
    for hx,hw,hl in zip(func_xs,func_ws,func_hdrs):
        _box(sl,hx,Inches(1.88),hw,Inches(0.28),SURFACE2)
        _tx(sl,hx+Pt(2),Inches(1.91),hw,Inches(0.24),hl,7,bold=True,color=ACCENT2)
    for i,fr in enumerate(s["func_rank"]):
        ry=Inches(2.19)+i*Inches(0.52)
        is_best=(i==0)
        rbg=RGBColor(0x07,0x16,0x0e) if is_best else (SURFACE if i%2==0 else BG)
        rborder=RGBColor(0x2a,0x6e,0x4a) if is_best else SURFACE2
        _box(sl,Inches(6.2),ry,Inches(7.1),Inches(0.49),rbg,rborder,1.5 if is_best else 0.5)
        fs=fr["avg_functional_score"]
        fc=GREEN if is_best else (YELLOW if fs>=0.45 else RED)
        ms=fr["avg_manual_score"]
        mc2=GREEN if ms>=0.82 else (YELLOW if ms>=0.70 else RED)
        cq=fr["avg_cq_pass_rate"]
        cc=GREEN if cq>=0.18 else (YELLOW if cq>=0.10 else RED)
        vals=[fr["model"],fr["config"],f'{fr["avg_numb"]:.2f}',f'{ms:.3f}',f'{cq:.3f}',f'{fs:.3f}']
        vcols=[_rgb(mc2) if is_best else WHITE, MUTED, MUTED, _rgb(mc2), _rgb(cc), _rgb(fc)]
        bolds=[True,False,False,False,False,True]
        for hx,hw,val,vc,bold in zip(func_xs,func_ws,vals,vcols,bolds):
            _tx(sl,hx+Pt(2),ry+Pt(6),hw-Pt(4),Inches(0.4),val,9,bold=bold,color=vc)
    _callout(sl,Inches(0.5),Inches(5.35),Inches(12.3),Inches(0.55),
             "Key finding: Llama-2-70B (cq-only) wins functional eval (0.534) despite poor structural quality — because it produces "
             "near-minimal ontologies (avg. 0.79 extra elements). Automated metrics reward parsimony, not modeling correctness.",YELLOW)
    _callout(sl,Inches(0.5),Inches(6.0),Inches(12.3),Inches(0.52),
             "Evaluation gap: functional scores cannot detect ontological modeling errors (bad domain/range, incoherent axioms). "
             "Expert human evaluation is essential to validate real-world utility.",RED)

    # ─── SLIDE 11 — OPEN-SOURCE MODEL SELECTION ───────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"11","Open-Source Model Selection for Human Evaluation")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.3),
        "Structural and functional evaluations identify different open-source winners. Two configurations selected — one per quality dimension.",
        12,color=MUTED)
    # Track A card
    _box(sl,Inches(0.5),Inches(1.6),Inches(5.9),Inches(4.3),RGBColor(0x04,0x12,0x08),RGBColor(0x2a,0x6e,0x4a),2)
    _tx(sl,Inches(0.65),Inches(1.7),Inches(5.6),Inches(0.3),"Track A — Structural / Creativity Winner",12,bold=True,color=GREEN)
    _tx(sl,Inches(0.65),Inches(2.05),Inches(5.6),Inches(0.45),"Llama-3.1-8B  +  scenario-only",18,bold=True,color=WHITE)
    track_a_pts=["Structural score 0.595 — best open-source scenario-only",
                 "OOPS score 0.405 — highest open-source in this config",
                 "Parse rate 78.6% across 14 scenarios",
                 "Tests whether scenario text alone suffices for ODP generation",
                 "Richer ontological output; higher creativity",
                 "Trade-off: more superfluous additions than minimal configs"]
    for k,pt in enumerate(track_a_pts):
        _tx(sl,Inches(0.65),Inches(2.56)+k*Inches(0.38),Inches(5.6),Inches(0.35),f"• {pt}",10,color=WHITE)
    # Track B card
    _box(sl,Inches(6.6),Inches(1.6),Inches(6.2),Inches(4.3),RGBColor(0x12,0x10,0x04),RGBColor(0x80,0x70,0x20),2)
    _tx(sl,Inches(6.75),Inches(1.7),Inches(5.9),Inches(0.3),"Track B — Functional Winner",12,bold=True,color=YELLOW)
    _tx(sl,Inches(6.75),Inches(2.05),Inches(5.9),Inches(0.45),"Llama-2-70B  +  cq-only",18,bold=True,color=WHITE)
    track_b_pts=["Functional score 0.534 — best among all evaluated configurations",
                 "CQ pass rate 0.202 — highest open-source cq-only result",
                 "Parse rate 85.7% in cq-only configuration",
                 "Average 0.79 extra elements: extremely minimal output",
                 "Structural quality known to be poor (bad OWL modeling)",
                 "Human experts will assess functional vs. structural divergence"]
    for k,pt in enumerate(track_b_pts):
        _tx(sl,Inches(6.75),Inches(2.56)+k*Inches(0.38),Inches(5.9),Inches(0.35),f"• {pt}",10,color=WHITE)
    _callout(sl,Inches(0.5),Inches(6.05),Inches(12.3),Inches(0.52),
             "Scientific motivation: structural vs. functional metrics select DIFFERENT winners. Human evaluation will determine which "
             "dimension better predicts real-world utility for ontology engineers.",GREEN)
    _callout(sl,Inches(0.5),Inches(6.65),Inches(12.3),Inches(0.45),
             "Scope reduction: 2 open-source configs (from 5) makes human evaluation feasible while probing opposite ends of the prompting spectrum.",ACCENT)

    # ─── SLIDE 11b — HUMAN EVALUATION PLAN ───────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"12","Human Evaluation Plan")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.28),
        "Expert and student evaluation of selected LLM-generated ODPs. Blind design: evaluators do not know if an ODP was LLM-generated or human-authored.",
        12,color=MUTED)
    # ── Three form cards ──────────────────────────────────────────────────────
    form_data=[
        ("Pattern Authors",ACCENT,[
            "Author role: lead / co-author",
            "Intent match (1\u20135): models original pattern intent?",
            "Scope appropriateness (1\u20135): too broad / narrow / right?",
            "Reuse readiness (1\u20135): usable without major rework?",
            "Adoption threshold: start from draft or scratch? (Y/M/N)",
            "Time saved: 0% / 1\u201325% / 26\u201350% / 51\u201375% / 76\u2013100%",
            "Missing core elements: free text"
        ]),
        ("Ontology Experts",GREEN,[
            "Overall helpfulness (1\u20135)",
            "Reusability (1\u20135): reuse in related context?",
            "Clarity / Documentation (1\u20135)",
            "Adoption threshold: (Yes / Maybe / No + justification)",
            "Time saved: 0% \u2192 76\u2013100%",
            "Free text: observations, concerns, suggestions"
        ]),
        ("MSc / PhD Students",YELLOW,[
            "Time to understand (minutes)",
            "Documentation clarity (1\u20135)",
            "Naming clarity (1\u20135): intuitive names?",
            "Adoption threshold: (Yes / Maybe / No + justification)",
            "Time saved: 0% \u2192 76\u2013100%"
        ]),
    ]
    for i,(title,color,bullets) in enumerate(form_data):
        cx=Inches(0.5+i*4.2)
        _card(sl,cx,Inches(1.55),Inches(4.0),Inches(3.8),title,bullets,_rgb(color),_dim3(color))
    # ── Pool and notes ────────────────────────────────────────────────────────
    _tx(sl,Inches(0.5),Inches(5.55),Inches(12.3),Inches(0.28),"Participant Pool",13,bold=True,color=ACCENT2)
    pool_items=[
        ("Pattern authors","~2 (incl. 1 also student)",ACCENT2),
        ("Ontology experts","~5 (some also authors)",GREEN),
        ("MSc/PhD students","~2",YELLOW),
        ("ISWC workshop","~8 experts + students",PURPLE),
        ("Target","~38 expert + ~38 student evaluators",WHITE),
        ("Load","~3 ontologies/person; authors: ~2",MUTED),
    ]
    for i,(lbl,val,vc) in enumerate(pool_items):
        col=i%3; row=i//3
        px=Inches(0.5+col*4.2); py=Inches(5.88+row*0.45)
        _box(sl,px,py,Inches(4.0),Inches(0.4),SURFACE,SURFACE2,0.5)
        _tx(sl,px+Pt(5),py+Pt(6),Inches(1.6),Inches(0.32),lbl,9,bold=True,color=MUTED)
        _tx(sl,px+Inches(1.7),py+Pt(6),Inches(2.2),Inches(0.32),val,9,bold=True,color=_rgb(vc))
    _callout(sl,Inches(0.5),Inches(7.0),Inches(8.5),Inches(0.38),
             "Open question: trim very long CQ lists (~3 representative CQs) to avoid participant fatigue.",YELLOW)
    _callout(sl,Inches(9.15),Inches(7.0),Inches(3.65),Inches(0.38),
             "Platform: github.com/ebrahimnorouzi/odp-platform",ACCENT)

    # ─── SLIDE 12 — MATCHED NAMES ─────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"12","Matched Class & Property Names vs Ground Truth")
    _tx(sl,Inches(0.5),Inches(1.2),Inches(12.3),Inches(0.3),
        "Names in both generated ontology and ground truth (CamelCase split, case-folded). Low F1 ≠ zero conceptual overlap.",12,color=MUTED)
    _card(sl,Inches(0.5),Inches(1.6),Inches(3.8),Inches(1.75),"Matching Logic",
          ["URI local name extracted: :hasCause → 'has cause'",
           "CamelCase split: HasEffect → 'has effect'","Case-folded & whitespace-normalised",
           "Precision=|matched|/|generated|, Recall=|matched|/|GT|"],ACCENT2,ACCENT)
    # examples
    _tx(sl,Inches(4.5),Inches(1.6),Inches(8.1),Inches(0.28),"Example Matches (GPT-5.4 & Gemini, cq-only)",12,bold=True,color=ACCENT2)
    ex_rows=[("GPT-5.4/cq-only/2023-133-01","event","has effect, has outcome"),
             ("GPT-5.4/cq-only/2023-135-01","name usage","has name usage"),
             ("GPT-5.4/cq-only/2025-149-01","relational database, table, column","—"),
             ("Gemini/scenario-cq/2023-133-01","event","—"),
             ("Gemini/cq-only/2025-151-01","evaluation request","requested action, issued by")]
    ex_xs=[Inches(4.5),Inches(7.1),Inches(10.1)]; ex_ws=[Inches(2.55),Inches(2.95),Inches(2.75)]
    for hx,hw,hl in zip(ex_xs,ex_ws,["Model/Config/Scenario","Matched Classes","Matched Properties"]):
        _box(sl,hx,Inches(1.93),hw,Inches(0.28),SURFACE2); _tx(sl,hx+Pt(3),Inches(1.96),hw,Inches(0.24),hl,8,bold=True,color=ACCENT2)
    for i,(tag,cls,prop) in enumerate(ex_rows):
        ry=Inches(2.24)+i*Inches(0.53)
        rbg=SURFACE if i%2==0 else BG
        for hx,hw in zip(ex_xs,ex_ws): _box(sl,hx,ry,hw,Inches(0.5),rbg)
        for hx,hw,val,vc in zip(ex_xs,ex_ws,[tag,cls,prop],[MUTED,GREEN,ACCENT2]):
            _tx(sl,hx+Pt(3),ry+Pt(4),hw-Pt(6),Inches(0.44),val,9,color=_rgb(vc))
    # prop label chart
    _tx(sl,Inches(0.5),Inches(3.5),Inches(3.8),Inches(0.28),"Property Label Ratio by Model",12,bold=True,color=ACCENT2)
    for i,m in enumerate(models):
        if qual.empty or m not in qual.index or ps.loc[m,"rate"]==0: continue
        plr=qual.loc[m,"property_label_ratio"] if "property_label_ratio" in qual.columns else 0
        _hbar(sl,Inches(0.5),Inches(3.85)+i*Inches(0.5),Inches(3.5),Inches(0.4),
              plr,_m(m,"name",m),_fmt(plr),_parse_bar_color(plr))
    _callout(sl,Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.38),
             "class_label_ratio=1.000 for ALL models on ALL successful parses. The split metric reveals property annotation gaps.",YELLOW)

    # ─── SLIDE 13 — SUMMARY ───────────────────────────────────────────────────
    sl = _blank(prs); _bg(sl); _header(sl,"13","Summary & Best Performing Model")
    wp=s["winners"].get("best_parse",""); wh=s["winners"].get("best_halluc",""); wf=s["winners"].get("best_f1","")
    _box(sl,Inches(0.5),Inches(1.22),Inches(5.5),Inches(5.5),RGBColor(0x07,0x16,0x0e),RGBColor(0x2a,0x6e,0x4a),2)
    _tx(sl,Inches(0.65),Inches(1.32),Inches(5.2),Inches(0.38),"★  Best Overall Reliability",14,bold=True,color=GREEN)
    _tx(sl,Inches(0.65),Inches(1.72),Inches(5.2),Inches(0.42),_m(wp,"name",wp),20,bold=True,color=WHITE)
    parse_r=_pct(ps.loc[wp,"rate"]) if wp in ps.index else "—"
    best_pts=[f"Parse success: {parse_r} (highest across all configs)",
              "100% parse rate on EVERY configuration","Best CQ coverage and Combined F1",
              "Trade-off: highest hallucination (0.623)","Adds unsupported domain knowledge"]
    for k,pt in enumerate(best_pts):
        _tx(sl,Inches(0.65),Inches(2.18)+k*Inches(0.42),Inches(5.2),Inches(0.38),f"• {pt}",10,color=WHITE)
    _box(sl,Inches(0.5),Inches(4.55),Inches(2.62),Inches(2.1),RGBColor(0x10,0x10,0x06),RGBColor(0x80,0x70,0x20),2)
    _tx(sl,Inches(0.62),Inches(4.65),Inches(2.4),Inches(0.32),"Faithfulness",11,bold=True,color=YELLOW)
    _tx(sl,Inches(0.62),Inches(4.98),Inches(2.4),Inches(0.3),_m(wh,"name",wh),15,bold=True,color=WHITE)
    hlv=_fmt(qual.loc[wh,"hallucination_ratio"]) if wh in qual.index and "hallucination_ratio" in qual.columns else "—"
    for k,pt in enumerate([f"Hallucination: {hlv} (lowest)","Best open-source model","Recommended: scenario-cq-reasoning"]):
        _tx(sl,Inches(0.62),Inches(5.32)+k*Inches(0.38),Inches(2.4),Inches(0.35),f"• {pt}",9,color=WHITE)
    _box(sl,Inches(3.27),Inches(4.55),Inches(2.65),Inches(2.1),RGBColor(0x10,0x0d,0x06),RGBColor(0x80,0x50,0x20),2)
    _tx(sl,Inches(3.38),Inches(4.65),Inches(2.4),Inches(0.32),"Best GT Match",11,bold=True,color=ACCENT2)
    _tx(sl,Inches(3.38),Inches(4.98),Inches(2.4),Inches(0.3),_m(wf,"name",wf),15,bold=True,color=WHITE)
    f1v=_fmt(qual.loc[wf,"combined_f1"]) if wf in qual.index and "combined_f1" in qual.columns else "—"
    for k,pt in enumerate([f"Combined F1: {f1v} (best)","Closest to expert ontology vocab","High CQ coverage"]):
        _tx(sl,Inches(3.38),Inches(5.32)+k*Inches(0.38),Inches(2.4),Inches(0.35),f"• {pt}",9,color=WHITE)
    findings=[
        (GREEN, "F1","Model size ≠ performance. Llama-3.1-8B (8B) beats Llama-2-70B (70B) on parse reliability."),
        (GREEN, "F2","Gemini 3.1 Pro: 100% parse, best F1 & CQ coverage — but highest hallucination (0.623)."),
        (YELLOW,"F3","class_label_ratio=1.0 for all models. property_label_ratio exposes hidden gaps (Llama-2: 0.717)."),
        (ACCENT2,"F4","scenario-cq-reasoning: lowest hallucination (0.461 avg), best faithfulness-quality trade-off."),
        (RED,   "F5","BLOOMZ-7B1: 0% parse across all 5 configs. Instruction-tuned BLOOM cannot produce valid Turtle."),
        (MUTED, "F6","Combined F1 universally low (0.06–0.13): LLMs diverge in naming from expert ontologists."),
    ]
    for k,(fc,tag,text) in enumerate(findings):
        fy=Inches(1.22)+k*Inches(0.78)
        _box(sl,Inches(6.3),fy,Inches(0.5),Inches(0.65),_dim(fc),_dim3(fc),1)
        _tx(sl,Inches(6.3),fy+Pt(8),Inches(0.5),Inches(0.5),tag,10,bold=True,color=_rgb(fc),align=PP_ALIGN.CENTER)
        _tx(sl,Inches(6.88),fy+Pt(8),Inches(5.85),Inches(0.58),text,11,color=WHITE)
    _callout(sl,Inches(0.5),Inches(7.0),Inches(12.3),Inches(0.38),
             "Recommendations — Proprietary: Gemini + scenario-only (structural) / GPT + scenario-cq-constraints (functional). "
             "Open-source for human eval: Llama-3.1-8B + scenario-only (structural track) and Llama-2-70B + cq-only (functional track, 0.534).",GREEN)

    prs.save(PPTX_OUT)
    print(f"Saved: {PPTX_OUT}  ({len(prs.slides)} slides)")


# ═══════════════════════════════════════════════════════════════════════════════
# ENTRY POINT
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    p = argparse.ArgumentParser(description="Generate ODPGen presentation files.")
    p.add_argument("--html-only", action="store_true")
    p.add_argument("--pptx-only", action="store_true")
    args = p.parse_args()

    print("Loading data from repo …")
    stats = load_stats()
    print(f"  Models: {stats['n_models_all']}  |  Configs: {stats['n_configs']}  |  "
          f"Scenarios: {stats['n_scenarios']}  |  Runs: {stats['n_runs']}  |  Metrics: {stats['n_metrics']}")

    if not args.pptx_only:
        build_html(stats)
    if not args.html_only:
        build_pptx(stats)


if __name__ == "__main__":
    main()
